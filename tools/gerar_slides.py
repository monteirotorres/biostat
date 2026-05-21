"""Monta o PPTX de Bioestatística no estilo dos slides originais:
fundo com marca-d'água UFRJ·IBCCF, títulos serifados, fórmulas
renderizadas como imagens e figuras estáticas no estilo do site.
~100 slides cobrindo todo o conteúdo do curso."""

import hashlib
from pathlib import Path
from PIL import Image
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

BASE = Path(__file__).parent.parent
FIG = BASE / "assets" / "slides"
FORM = FIG / "formulas"
FORM.mkdir(parents=True, exist_ok=True)
FUNDO = FIG / "fundo.png"
LOGO = BASE / "assets" / "ibccf-logo.png"
OUT = BASE / "bioestatistica_slides.pptx"

BG = RGBColor(0xF3, 0xEC, 0xDF)
PAPER = RGBColor(0xFF, 0xFD, 0xF8)
INK = RGBColor(0x1A, 0x1A, 0x1A)
BLUE = RGBColor(0x32, 0x66, 0xAD)
RED = RGBColor(0xC0, 0x39, 0x2B)
MUTED = RGBColor(0x6B, 0x64, 0x57)
LINE = RGBColor(0xC9, 0xBF, 0xA9)
TINT = RGBColor(0xEC, 0xF1, 0xF8)   # azul bem claro p/ caixas de exemplo

SERIF = "Georgia"
MONO = "Consolas"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]
_num = 0


# ---------------- fórmulas como imagem ----------------
def formula_png(latex, fontsize=30, color="#1a1a1a"):
    key = hashlib.md5(f"{latex}{fontsize}{color}".encode()).hexdigest()[:12]
    path = FORM / f"f_{key}.png"
    if not path.exists():
        fig = plt.figure(figsize=(0.01, 0.01))
        fig.text(0, 0, f"${latex}$", fontsize=fontsize, color=color)
        fig.savefig(path, dpi=200, transparent=True, bbox_inches="tight", pad_inches=0.06)
        plt.close(fig)
    return path


# ---------------- helpers de slide ----------------
def _bg(slide):
    slide.shapes.add_picture(str(FUNDO), 0, 0, SW, SH)
    pic = slide.shapes[-1]
    slide.shapes._spTree.remove(pic._element)
    slide.shapes._spTree.insert(2, pic._element)


def _box(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    return tb


def _set(p, text, size, color=INK, bold=False, font=SERIF, italic=False):
    p.text = text
    for r in p.runs:
        r.font.size = Pt(size); r.font.color.rgb = color
        r.font.bold = bold; r.font.name = font; r.font.italic = italic


def _rect(slide, l, t, w, h, color, rounded=False):
    shp = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    sh = slide.shapes.add_shape(shp, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background(); sh.shadow.inherit = False
    return sh


def _img_fit(slide, path, l, t, w, h):
    iw, ih = Image.open(path).size
    ar = iw / ih; box_ar = w / h
    if ar > box_ar:
        nw = w; nh = int(w / ar)
    else:
        nh = h; nw = int(h * ar)
    slide.shapes.add_picture(str(path), l + (w - nw) // 2, t + (h - nh) // 2, nw, nh)
    return nw, nh


def _footer(slide):
    global _num
    _num += 1
    tb = _box(slide, Inches(0.5), Inches(7.04), Inches(6), Inches(0.4))
    _set(tb.text_frame.paragraphs[0], "Bioestatística · IBCCF · UFRJ", 10, MUTED, font=MONO)
    nb = _box(slide, Inches(12.3), Inches(7.04), Inches(0.8), Inches(0.4))
    p = nb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    _set(p, str(_num), 10, MUTED, font=MONO)


def _bullets(slide, items, l, t, w, h, size=18):
    tb = _box(slide, l, t, w, h)
    tf = tb.text_frame
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        sub = it.startswith("- ")
        txt = it[2:] if sub else it
        _set(p, ("–  " if sub else "•  ") + txt, size - (2 if sub else 0), INK)
        p.space_after = Pt(9); p.line_spacing = 1.12
        if sub:
            p.level = 1


def _formula_card(slide, latex, cx_l, cx_w, top, target_h_in=0.7, fontsize=30):
    path = formula_png(latex, fontsize=fontsize)
    iw, ih = Image.open(path).size
    h = Inches(target_h_in); w = int(h * iw / ih)
    maxw = cx_w - Inches(0.8)
    if w > maxw:
        w = maxw; h = int(w * ih / iw)
    card_w = w + Inches(0.7); card_h = h + Inches(0.45)
    card_l = cx_l + (cx_w - card_w) // 2
    _rect(slide, card_l, top, card_w, card_h, PAPER, rounded=True)
    slide.shapes.add_picture(str(path), card_l + (card_w - w) // 2, top + (card_h - h) // 2, w, h)
    return card_h


def _example(slide, text, l, t, w):
    card = _rect(slide, l, t, w, Inches(1.15), TINT, rounded=True)
    tb = _box(slide, l + Inches(0.2), t + Inches(0.08), w - Inches(0.4), Inches(1.0))
    tf = tb.text_frame
    _set(tf.paragraphs[0], "EXEMPLO", 11, BLUE, bold=True, font=MONO)
    p = tf.add_paragraph(); _set(p, text, 14, INK); p.line_spacing = 1.1


# ---------------- tipos de slide ----------------
def title_slide():
    s = prs.slides.add_slide(BLANK); _bg(s)
    if LOGO.exists():
        _img_fit(s, LOGO, Inches(4.67), Inches(1.0), Inches(4.0), Inches(1.25))
    tb = _box(s, Inches(1), Inches(2.7), Inches(11.33), Inches(1.6))
    _set(tb.text_frame.paragraphs[0], "Bioestatística", 60, INK, bold=True)
    tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    sub = _box(s, Inches(1), Inches(4.2), Inches(11.33), Inches(0.6))
    p = sub.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    _set(p, "Instituto de Biofísica Carlos Chagas Filho · UFRJ", 18, MUTED, font=MONO)
    au = _box(s, Inches(1), Inches(5.2), Inches(11.33), Inches(0.8))
    p = au.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    _set(p, "Pedro Torres   ·   Gilberto Weissmuller", 18, INK)


def section_slide(num, title):
    s = prs.slides.add_slide(BLANK); _bg(s)
    _rect(s, Inches(1), Inches(3.05), Inches(1.5), Inches(0.12), BLUE)
    nb = _box(s, Inches(1), Inches(2.05), Inches(4), Inches(1.0))
    _set(nb.text_frame.paragraphs[0], f"Parte {num}", 22, BLUE, font=MONO)
    tb = _box(s, Inches(1), Inches(3.35), Inches(11.3), Inches(1.6))
    _set(tb.text_frame.paragraphs[0], title, 46, INK, bold=True)
    _footer(s)


def content(eyebrow, title, bullets=None, formula=None, image=None, caption=None, example=None):
    s = prs.slides.add_slide(BLANK); _bg(s)
    eb = _box(s, Inches(0.7), Inches(0.42), Inches(11), Inches(0.4))
    _set(eb.text_frame.paragraphs[0], eyebrow.upper(), 12, MUTED, font=MONO)
    tb = _box(s, Inches(0.7), Inches(0.78), Inches(12), Inches(1.0))
    _set(tb.text_frame.paragraphs[0], title, 28, INK, bold=True)
    _rect(s, Inches(0.72), Inches(1.7), Inches(1.1), Inches(0.05), BLUE)
    top = Inches(2.05)

    if image and bullets:
        _bullets(s, bullets, Inches(0.7), top, Inches(5.5), Inches(4.6))
        _img_fit(s, FIG / image, Inches(6.5), top, Inches(6.3), Inches(4.4))
    elif image:
        _img_fit(s, FIG / image, Inches(1.3), top, Inches(10.7), Inches(4.4))
        if caption:
            cb = _box(s, Inches(1), Inches(6.5), Inches(11.33), Inches(0.4))
            p = cb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            _set(p, caption, 12, MUTED, italic=True)
    elif bullets and formula:
        _bullets(s, bullets, Inches(0.9), top, Inches(11.5), Inches(3.0), size=19)
        y = top + Inches(3.0)
        _formula_card(s, formula, Inches(0.9), Inches(11.5), y, target_h_in=0.75)
    elif formula:
        _formula_card(s, formula, Inches(1.5), Inches(10.3), Inches(3.0), target_h_in=1.0, fontsize=34)
    elif bullets:
        _bullets(s, bullets, Inches(0.9), top, Inches(11.5), Inches(4.6), size=20)

    if example:
        _example(s, example, Inches(0.9), Inches(5.85), Inches(11.5))
    _footer(s)


def refs_slide():
    s = prs.slides.add_slide(BLANK); _bg(s)
    tb = _box(s, Inches(0.7), Inches(0.7), Inches(11), Inches(1))
    _set(tb.text_frame.paragraphs[0], "Referências e contato", 28, INK, bold=True)
    _rect(s, Inches(0.72), Inches(1.6), Inches(1.1), Inches(0.05), BLUE)
    if LOGO.exists():
        _img_fit(s, LOGO, Inches(0.9), Inches(2.3), Inches(4.5), Inches(1.5))
    cb = _box(s, Inches(0.9), Inches(4.1), Inches(11), Inches(2.4))
    tf = cb.text_frame
    _set(tf.paragraphs[0], "Material teórico completo e interativo:", 18, INK)
    p = tf.add_paragraph(); _set(p, "monteirotorres.github.io/biostat", 18, BLUE, font=MONO)
    p = tf.add_paragraph(); _set(p, " ", 8)
    p = tf.add_paragraph(); _set(p, "Notebooks em Python (pandas · seaborn · scipy.stats)", 15, MUTED)
    p = tf.add_paragraph(); _set(p, "Contato: monteirotorres@biof.ufrj.br", 15, MUTED, font=MONO)
    _footer(s)


# =====================================================================
#                              CONTEÚDO
# =====================================================================
title_slide()

# ---------------------------- PARTE 1 ----------------------------
section_slide(1, "Introdução")

content("1 · Introdução", "Objeto da estatística", bullets=[
    "A estatística trata de coletar, resumir e interpretar dados na presença de variabilidade.",
    "Quase nunca medimos todos os indivíduos de interesse — trabalhamos com amostras.",
    "Duas grandes tarefas: descrever os dados e inferir propriedades da população.",
    "Toda conclusão carrega incerteza, que a estatística quantifica.",
])
content("1 · Introdução", "Amostra e população", bullets=[
    "População: o conjunto completo de indivíduos de interesse.",
    "Amostra: o subconjunto que efetivamente medimos.",
    "Medir toda a população costuma ser caro, demorado ou impossível.",
    "A amostragem aleatória é o que torna a generalização válida.",
], example="Para estimar a altura média dos adultos do país, medimos alguns milhares de pessoas (amostra) e inferimos a média de todos (população).")
content("1 · Introdução", "Parâmetro vs. estatística", bullets=[
    "Parâmetro: valor verdadeiro da população, em geral desconhecido.",
    "Estatística: valor calculado a partir da amostra, usado para estimar o parâmetro.",
    "A média da população é μ; a da amostra é x̄. O desvio padrão é σ (população) e s (amostra).",
], formula=r"\mu,\ \sigma \ \longrightarrow\ \bar{x},\ s")
content("1 · Introdução", "Descrição vs. inferência", bullets=[
    "Estatística descritiva: resume os dados em mãos (médias, gráficos, tabelas).",
    "Estatística inferencial: generaliza da amostra para a população.",
    "A descritiva fala só dos dados coletados; a inferencial vai além deles.",
    "Regra de ouro: descreva os dados antes de inferir qualquer coisa.",
])
content("1 · Introdução", "Precisão e acurácia", image="precisao.png",
        caption="Acurácia = proximidade do valor verdadeiro (viés). Precisão = consistência entre medidas (ruído).")
content("1 · Introdução", "Precisão e acurácia — em números", bullets=[
    "Acurácia se relaciona ao viés: a diferença entre a média das medidas e o valor verdadeiro.",
    "Precisão se relaciona à dispersão: o desvio padrão das medidas.",
    "Mais medidas reduzem o erro aleatório, mas não corrigem um viés sistemático.",
], formula=r"\mathrm{vi\acute{e}s} = \bar{x} - \mu \qquad \mathrm{precis\~ao} \sim s")
content("1 · Introdução", "Algarismos significativos", bullets=[
    "Refletem a precisão real do instrumento; dígitos a mais são ruído.",
    "Zeros à esquerda não contam; zeros à direita após a vírgula contam.",
    "Soma/subtração segue as casas decimais; multiplicação/divisão segue os significativos.",
    "Notação científica elimina a ambiguidade dos zeros.",
], example="0,00025 g tem 2 algarismos significativos; em microgramas (250 µg) continua tendo 2.")
content("1 · Introdução", "Notação científica", bullets=[
    "Escreve números muito grandes ou pequenos como a × 10ⁿ, com 1 ≤ |a| < 10.",
    "Expoente positivo: número grande; negativo: número pequeno.",
    "Em multiplicação e divisão, os expoentes somam e subtraem.",
], formula=r"6{,}022\times10^{23}\qquad 1{,}6\times10^{-19}")
content("1 · Introdução", "Tipos de variáveis", bullets=[
    "Qualitativas nominais: categorias sem ordem (tipo sanguíneo).",
    "Qualitativas ordinais: categorias com ordem (grau de dor).",
    "Quantitativas discretas: contagens (nº de células).",
    "Quantitativas contínuas: medidas (peso, pressão).",
    "O tipo determina o gráfico e o teste adequados.",
])
content("1 · Introdução", "Histogramas", image="histograma.png",
        caption="Revelam forma, centro, dispersão e outliers. Variável contínua, barras encostadas.")
content("1 · Introdução", "Curva normal", image="normal.png",
        caption="Definida por μ e σ; simétrica; vale a regra 68–95–99,7.")
content("1 · Introdução", "Curva normal — padronização", bullets=[
    "Qualquer normal pode ser convertida na normal padrão (μ = 0, σ = 1).",
    "O escore Z diz a quantos desvios padrão um valor está da média.",
    "Permite comparar valores de distribuições diferentes e consultar tabelas.",
], formula=r"Z = \frac{X - \mu}{\sigma}")
content("1 · Introdução", "Técnicas de amostragem", bullets=[
    "Aleatória simples: todos têm a mesma chance.",
    "Estratificada: sorteio dentro de subgrupos homogêneos.",
    "Conglomerados: sorteiam-se grupos inteiros.",
    "Aleatorização neutraliza variáveis de confusão; evite voluntários.",
])
content("1 · Introdução", "Variáveis de confusão", bullets=[
    "Afetam o desfecho e, ao mesmo tempo, diferem entre os grupos comparados.",
    "Criam associações enganosas e impedem conclusões causais.",
    "A aleatorização as distribui igualmente, em média, entre os grupos.",
], example="Cobradores de ônibus tinham menos doença cardíaca que motoristas — mas idade, estresse e saúde prévia são confundidores; o estudo mostra associação, não causa.")

# ---------------------------- PARTE 2 ----------------------------
section_slide(2, "Estatística Descritiva")

content("2 · Descritiva", "Medidas de tendência central", image="tendencia.png",
        caption="Moda, média e mediana resumem o centro de formas diferentes.")
content("2 · Descritiva", "Moda", bullets=[
    "É o valor que aparece com maior frequência.",
    "Pode ser única, múltipla (bimodal) ou inexistente.",
    "É a única medida de centro que funciona em variáveis categóricas.",
    "No histograma, corresponde ao pico da distribuição.",
])
content("2 · Descritiva", "Média aritmética", bullets=[
    "Soma de todos os valores dividida pela quantidade.",
    "Usa todos os dados, mas é sensível a valores extremos.",
    "É o ponto de equilíbrio: a soma dos desvios em relação a ela é zero.",
], formula=r"\bar{x} = \frac{1}{n}\sum_{i=1}^{n} x_i")
content("2 · Descritiva", "Média ponderada", bullets=[
    "Cada valor entra com um peso que reflete sua importância.",
    "Quando os pesos são iguais, recai na média aritmética.",
    "Base do valor esperado quando os pesos são probabilidades.",
], formula=r"\bar{x}_w = \frac{\sum w_i\,x_i}{\sum w_i}")
content("2 · Descritiva", "Média geométrica", bullets=[
    "Apropriada para dados multiplicativos: taxas de crescimento, diluições, títulos.",
    "É a média aritmética dos logaritmos, revertida pela exponencial.",
    "Sempre ≤ média aritmética; só vale para valores positivos.",
], formula=r"\bar{x}_g = \sqrt[n]{x_1 x_2 \cdots x_n}")
content("2 · Descritiva", "Média cortada", bullets=[
    "Remove uma porcentagem dos valores extremos antes de calcular a média.",
    "Mais robusta a outliers que a média comum.",
    "Cortar 0% é a média; cortar 50% é a mediana.",
], example="Em provas de ginástica artística, descartam-se a maior e a menor nota dos juízes.")
content("2 · Descritiva", "Mediana", bullets=[
    "Valor central com os dados ordenados; 50% abaixo e 50% acima.",
    "Com n par, é a média dos dois valores centrais.",
    "Robusta a outliers — pouco afetada por valores extremos.",
])
content("2 · Descritiva", "Média x mediana", bullets=[
    "Distribuição simétrica: média ≈ mediana.",
    "Assimétrica à direita: média > mediana (renda, tempo de espera).",
    "Assimétrica à esquerda: média < mediana.",
    "Sempre olhe o histograma antes de escolher.",
])
content("2 · Descritiva", "Quartis e percentis", bullets=[
    "Quartis dividem os dados ordenados em quatro partes iguais.",
    "Q1 (percentil 25), Q2 (mediana), Q3 (percentil 75).",
    "São a base do boxplot e descrevem a distribuição sem supor normalidade.",
])
content("2 · Descritiva", "Teorema do limite central", image="tlc.png",
        caption="A média de amostras tende à normal, mesmo com população não normal.")
content("2 · Descritiva", "TLC — erro padrão da média", bullets=[
    "A distribuição das médias tem a mesma média μ da população.",
    "Seu desvio padrão é o erro padrão, que diminui com a raiz de n.",
    "Por isso amostras maiores dão estimativas mais precisas.",
], formula=r"\bar{x} \sim \mathcal{N}\!\left(\mu,\ \frac{\sigma^2}{n}\right)\qquad \mathrm{EP} = \frac{\sigma}{\sqrt{n}}")
content("2 · Descritiva", "Amplitude", bullets=[
    "Maior valor menos o menor valor.",
    "Simples, mas usa apenas dois valores e é muito sensível a outliers.",
    "Cresce com o tamanho da amostra.",
], formula=r"A = x_{\max} - x_{\min}")
content("2 · Descritiva", "Desvio médio", bullets=[
    "Média das distâncias absolutas até a média.",
    "Intuitivo, mas pouco usado por causa do valor absoluto.",
    "Levou ao uso do quadrado dos desvios — variância e desvio padrão.",
], formula=r"\mathrm{DM} = \frac{1}{n}\sum |x_i - \bar{x}|")
content("2 · Descritiva", "Variância", bullets=[
    "Média dos quadrados dos desvios em relação à média.",
    "Penaliza desvios grandes; fica numa unidade ao quadrado.",
    "A amostral divide por n − 1 (correção de Bessel).",
], formula=r"s^2 = \frac{1}{n-1}\sum (x_i - \bar{x})^2")
content("2 · Descritiva", "Desvio padrão", bullets=[
    "Raiz da variância — volta à unidade original dos dados.",
    "Na normal, vale a regra 68–95–99,7.",
    "É a medida de dispersão mais reportada em ciência.",
], formula=r"s = \sqrt{\frac{1}{n-1}\sum (x_i - \bar{x})^2}")
content("2 · Descritiva", "Desvio padrão e erro padrão", bullets=[
    "Desvio padrão (s): variabilidade dos indivíduos.",
    "Erro padrão da média (EP): variabilidade da média amostral.",
    "O EP é sempre menor e diminui com amostras maiores.",
    "Em gráficos, sempre diga se a barra é desvio padrão, EP ou IC.",
], formula=r"\mathrm{EP} = \frac{s}{\sqrt{n}}")
content("2 · Descritiva", "Boxplot e desvio interquartil", image="boxplot.png",
        caption="IQR = Q3 − Q1 descreve os 50% centrais e é robusto a outliers.")
content("2 · Descritiva", "Coeficiente de variação", bullets=[
    "Desvio padrão dividido pela média, em porcentagem — dispersão relativa.",
    "Compara variabilidade entre grupos de escalas diferentes.",
    "Só faz sentido em variáveis com zero absoluto.",
], formula=r"\mathrm{CV} = \frac{s}{\bar{x}}\times 100\%")

# ---------------------------- PARTE 3 ----------------------------
section_slide(3, "Probabilidade")

content("3 · Probabilidade", "Evento e espaço amostral", bullets=[
    "Experimento aleatório: resultado não previsível, mas com resultados conhecidos.",
    "Espaço amostral (Ω): conjunto de todos os resultados possíveis.",
    "Evento: um subconjunto do espaço amostral.",
    "A probabilidade é um número entre 0 (impossível) e 1 (certo).",
])
content("3 · Probabilidade", "Regras de probabilidade", bullets=[
    "Complemento: a chance de não ocorrer é 1 menos a chance de ocorrer.",
    "União: soma das probabilidades menos a interseção.",
    "Eventos independentes: a interseção é o produto das probabilidades.",
], formula=r"P(A\cup B) = P(A)+P(B)-P(A\cap B)")
content("3 · Probabilidade", "Probabilidade condicional", bullets=[
    "Mede a chance de A sabendo que B ocorreu.",
    "Restringe o espaço amostral ao subconjunto em que B é verdade.",
    "Cuidado: P(A | B) geralmente é diferente de P(B | A).",
], formula=r"P(A\mid B) = \frac{P(A\cap B)}{P(B)}")
content("3 · Probabilidade", "Teorema de Bayes", bullets=[
    "Inverte a direção de uma probabilidade condicional.",
    "Permite atualizar a probabilidade de uma hipótese diante de uma evidência.",
    "É a base do raciocínio diagnóstico.",
], formula=r"P(A\mid B) = \frac{P(B\mid A)\,P(A)}{P(B)}")
content("3 · Probabilidade", "Análise combinatória", bullets=[
    "Permutações: contam arranjos em que a ordem importa.",
    "Combinações: contam escolhas em que a ordem não importa.",
    "Base do cálculo de probabilidades na distribuição binomial.",
], formula=r"\binom{n}{k} = \frac{n!}{k!\,(n-k)!}", example="Mega-Sena: C(60,6) = 50 063 860 resultados; a chance de acertar é 1 nesse total.")
content("3 · Probabilidade", "Testes binários", bullets=[
    "Resultado com apenas dois valores: positivo ou negativo.",
    "Quatro cenários: verdadeiro/falso positivo e verdadeiro/falso negativo.",
    "Reduzir um tipo de erro tipicamente aumenta o outro.",
])
content("3 · Probabilidade", "Sensibilidade e especificidade", image="sens_esp.png",
        caption="O ponto de corte equilibra falsos positivos e falsos negativos.")
content("3 · Probabilidade", "Sensibilidade e especificidade — fórmulas", bullets=[
    "Sensibilidade: proporção de doentes que o teste detecta.",
    "Especificidade: proporção de saudáveis que o teste descarta.",
    "São propriedades do teste, independentes da prevalência.",
], formula=r"\mathrm{Sens} = \frac{VP}{VP+FN}\qquad \mathrm{Esp} = \frac{VN}{VN+FP}")
content("3 · Probabilidade", "Matriz de confusão", bullets=[
    "Tabela 2×2 que mostra acertos e erros do teste.",
    "Diagonal = acertos; fora dela, onde o teste se confundiu.",
    "Dela derivam acurácia, sensibilidade, especificidade, VPP e VPN.",
    "Acurácia sozinha engana quando as classes são desbalanceadas.",
])
content("3 · Probabilidade", "Valor preditivo", bullets=[
    "VPP: probabilidade de ser doente dado um teste positivo.",
    "VPN: probabilidade de ser saudável dado um teste negativo.",
    "Dependem fortemente da prevalência da doença.",
], formula=r"\mathrm{VPP} = \frac{VP}{VP+FP}")
content("3 · Probabilidade", "O paradoxo da doença rara", bullets=[
    "Mesmo um teste excelente gera muitos falsos positivos quando a doença é rara.",
    "A maioria dos positivos vem da enorme população de saudáveis.",
    "Por isso testes de triagem são aplicados em grupos de maior risco.",
], example="Sensibilidade e especificidade de 99%, prevalência de 0,1%: apenas ~9% dos positivos são realmente doentes.")
content("3 · Probabilidade", "Curva ROC", image="roc.png",
        caption="Desempenho do teste em todos os cortes; a AUC resume a qualidade.")
content("3 · Probabilidade", "Distribuição binomial", image="binomial.png",
        caption="Número de sucessos em n tentativas independentes com probabilidade p.")
content("3 · Probabilidade", "Binomial — fórmula", bullets=[
    "Conta k sucessos em n tentativas independentes de mesma probabilidade p.",
    "A média é np e a variância é np(1 − p).",
    "Para n grande, aproxima-se da normal.",
], formula=r"P(X=k) = \binom{n}{k} p^{k}(1-p)^{n-k}")

# ---------------------------- PARTE 4 ----------------------------
section_slide(4, "Distribuições amostrais")

content("4 · Distribuições amostrais", "De onde vêm t, χ² e F", image="dist_amostral.png",
        caption="Todas surgem da amostragem de variáveis normais.")
content("4 · Distribuições amostrais", "Graus de liberdade", bullets=[
    "Número de valores livres ao calcular uma estatística.",
    "Com n valores e a média fixada, só n − 1 desvios são livres.",
    "Quanto maiores os graus de liberdade, mais as distribuições se aproximam da normal.",
])
content("4 · Distribuições amostrais", "Distribuição qui-quadrado", bullets=[
    "Soma de quadrados de variáveis normais padrão.",
    "Só assume valores positivos; assimétrica à direita.",
    "Aparece em variâncias e no teste qui-quadrado.",
], formula=r"\chi^2_k = Z_1^2 + Z_2^2 + \cdots + Z_k^2")
content("4 · Distribuições amostrais", "Distribuição t de Student", bullets=[
    "Surge ao padronizar a média usando o desvio padrão estimado s.",
    "Tem caudas mais pesadas que a normal — exige mais evidência.",
    "Com muitos graus de liberdade, converge para a normal.",
], formula=r"t = \frac{\bar{x}-\mu}{s/\sqrt{n}}")
content("4 · Distribuições amostrais", "Distribuição F", bullets=[
    "Razão entre duas variâncias (duas qui-quadrado sobre seus gl).",
    "Base da ANOVA e da comparação de variâncias.",
    "Próxima de 1 quando as variâncias estimam a mesma coisa.",
], formula=r"F = \frac{\chi^2_{d_1}/d_1}{\chi^2_{d_2}/d_2}")

# ---------------------------- PARTE 5 ----------------------------
section_slide(5, "Estatística Inferencial")

content("5 · Inferencial", "Lógica da inferência", bullets=[
    "Assume-se H₀ verdadeira e calcula-se a probabilidade dos dados observados.",
    "Se essa probabilidade é muito baixa, rejeita-se H₀.",
    "É a lógica do júri: inocente até que as evidências afastem a dúvida.",
    "É um filtro consensual contra o acaso e o viés.",
])
content("5 · Inferencial", "Estimação", bullets=[
    "Pontual: um único número estima o parâmetro (x̄ estima μ).",
    "Intervalar: um intervalo de valores plausíveis.",
    "Um bom estimador é não viesado, eficiente e consistente.",
])
content("5 · Inferencial", "Intervalo de confiança", image="ic.png",
        caption="IC 95%: ao repetir o estudo, ~95% dos intervalos conteriam o parâmetro.")
content("5 · Inferencial", "Intervalo de confiança — fórmula", bullets=[
    "Estimativa pontual mais ou menos a margem de erro.",
    "A margem depende da variabilidade, do tamanho da amostra e do nível de confiança.",
    "Quadruplicar n reduz a margem pela metade.",
], formula=r"\bar{x} \pm t^{*}\,\frac{s}{\sqrt{n}}")
content("5 · Inferencial", "Hipótese nula e alternativa", bullets=[
    "H₀: ausência de efeito ou diferença (o status quo).",
    "H₁: existe efeito ou diferença.",
    "H₀ é sempre uma igualdade; H₁ pode ser bilateral ou unilateral.",
    "Não rejeitar H₀ não é o mesmo que provar H₀.",
])
content("5 · Inferencial", "Nível de significância (α)", bullets=[
    "É o risco aceito de rejeitar H₀ quando ela é verdadeira (erro tipo I).",
    "Valor mais comum: 0,05. Áreas exigentes usam 0,01 ou menos.",
    "Deve ser fixado antes de coletar os dados.",
    "Relaciona-se ao IC: confiança + α = 1.",
])
content("5 · Inferencial", "Valor-p", image="valorp.png",
        caption="Probabilidade de dados tão ou mais extremos que os observados, sob H₀.")
content("5 · Inferencial", "Valor-p — o que NÃO é", bullets=[
    "Não é a probabilidade de H₀ ser verdadeira.",
    "Não é a probabilidade de o resultado ser por acaso.",
    "Não mede o tamanho nem a importância do efeito.",
    "Reporte sempre o tamanho de efeito e o intervalo de confiança.",
])
content("5 · Inferencial", "Erros tipo I e II", image="erros.png",
        caption="α = falso positivo; β = falso negativo; poder = 1 − β.")
content("5 · Inferencial", "Poder estatístico", bullets=[
    "Probabilidade de detectar um efeito que realmente existe.",
    "Cresce com o tamanho da amostra, o tamanho do efeito e a redução da variabilidade.",
    "Meta usual: poder ≥ 80%.",
], formula=r"\text{poder} = 1 - \beta")
content("5 · Inferencial", "Teste z", bullets=[
    "Compara uma média com um valor de referência quando σ é conhecido.",
    "A estatística segue a normal padrão sob H₀.",
    "Raro em pesquisa, pois σ quase nunca é conhecido.",
], formula=r"z = \frac{\bar{x}-\mu_0}{\sigma/\sqrt{n}}")
content("5 · Inferencial", "Teste t de uma amostra", bullets=[
    "Compara a média de uma amostra com um valor de referência.",
    "Usa o desvio padrão amostral s; segue a distribuição t.",
    "É robusto a desvios moderados da normalidade.",
], formula=r"t = \frac{\bar{x}-\mu_0}{s/\sqrt{n}}")
content("5 · Inferencial", "Teste t para dois grupos", bullets=[
    "Compara as médias de dois grupos independentes.",
    "A versão de Welch não supõe variâncias iguais — recomendada por padrão.",
    "Reporte também o tamanho de efeito (d de Cohen).",
], formula=r"t = \frac{\bar{x}_1-\bar{x}_2}{\sqrt{s_1^2/n_1 + s_2^2/n_2}}")
content("5 · Inferencial", "Teste t pareado", bullets=[
    "Cada indivíduo é medido duas vezes (antes e depois).",
    "Testa se a média das diferenças é zero.",
    "Elimina a variabilidade entre indivíduos — mais potente.",
], example="Pressão arterial antes e depois de um remédio, no mesmo paciente.")
content("5 · Inferencial", "ANOVA", image="anova.png",
        caption="Compara 3+ grupos sem inflar o erro tipo I.")
content("5 · Inferencial", "ANOVA — a estatística F", bullets=[
    "Compara a variância entre grupos com a variância dentro dos grupos.",
    "Se as médias são iguais, F ≈ 1; se diferem, F cresce.",
    "Rejeitar H₀ indica que alguma média difere — testes post hoc dizem qual.",
], formula=r"F = \frac{s^2_{\mathrm{entre}}}{s^2_{\mathrm{dentro}}}")
content("5 · Inferencial", "Teste do qui-quadrado", bullets=[
    "Testa associação entre variáveis categóricas.",
    "Compara contagens observadas com as esperadas sob independência.",
    "Depende do número absoluto de observações, não só das proporções.",
], formula=r"\chi^2 = \sum \frac{(O_i - E_i)^2}{E_i}")
content("5 · Inferencial", "Correlação", image="correlacao.png",
        caption="r de Pearson mede a associação linear; resíduos em vermelho.")
content("5 · Inferencial", "Correlação — coeficiente de Pearson", bullets=[
    "Varia de −1 a +1; o sinal indica a direção da associação.",
    "r² é a fração da variância de uma variável explicada pela outra.",
    "Correlação não implica causalidade.",
    "Pearson só capta relações lineares — sempre faça o gráfico.",
], formula=r"-1 \leq r \leq 1")
content("5 · Inferencial", "Regressão linear", image="regressao.png",
        caption="Ajusta uma reta por mínimos quadrados; R² mede a qualidade do ajuste.")
content("5 · Inferencial", "Regressão — mínimos quadrados", bullets=[
    "A reta minimiza a soma dos quadrados dos resíduos.",
    "A inclinação diz quanto y muda por unidade de x.",
    "Não extrapole muito além do intervalo observado.",
], formula=r"\hat{y} = \hat{\beta}_0 + \hat{\beta}_1 x")
content("5 · Inferencial", "Testes não paramétricos", bullets=[
    "Dispensam a suposição de normalidade; trabalham com postos.",
    "Mann–Whitney substitui o t independente; Wilcoxon, o pareado.",
    "Kruskal–Wallis substitui a ANOVA.",
    "Menos potentes quando a normalidade vale, mas mais robustos.",
])
content("5 · Inferencial", "Design experimental", bullets=[
    "Pilares: controle, aleatorização, repetição e pareamento.",
    "Experimento verdadeiro permite causa; estudo observacional, só associação.",
    "Significância estatística não é o mesmo que relevância prática.",
])
content("5 · Inferencial", "Cálculo amostral", bullets=[
    "Define quantos indivíduos são necessários para um dado poder.",
    "Depende do tamanho de efeito, da variabilidade e de α.",
    "Deve ser feito antes do estudo — estudos subdimensionados são problemáticos.",
], formula=r"n \approx 2\left(\frac{z_{1-\alpha/2}+z_{1-\beta}}{\delta/\sigma}\right)^2")

# ---------------------------- TABELAS ----------------------------
section_slide(6, "Tabelas estatísticas")
content("6 · Tabelas", "Valores críticos mais usados", bullets=[
    "z para IC 90% = 1,645; 95% = 1,960; 99% = 2,576.",
    "Tabela t: valores maiores que z para amostras pequenas, convergindo a z.",
    "Tabela χ²: cauda direita, usada em aderência e contingência.",
    "Tabela F: usada em ANOVA, com dois graus de liberdade.",
])
content("6 · Tabelas", "Tabela ou computador?", bullets=[
    "Hoje usamos o computador (scipy.stats) para qualquer valor crítico.",
    "Saber ler uma tabela ajuda a entender o que o software faz.",
    "stats.norm.ppf, stats.t.ppf, stats.chi2.ppf e stats.f.ppf dão os valores.",
], formula=r"P(Z \leq 1{,}96) = 0{,}975")

# ---------------------------- EXERCÍCIOS ----------------------------
section_slide(7, "Exercícios resolvidos")

content("7 · Exercícios", "Algarismos significativos", bullets=[
    "Numa balança de 5 casas lê-se 0,00025 g.",
    "Quantos algarismos significativos tem a leitura? E em microgramas (250 µg)?",
], example="2 algarismos significativos. Trocar a unidade não muda nada: 250 µg também tem 2 — é uma propriedade da medida, não da unidade.")
content("7 · Exercícios", "Probabilidade condicional (Bayes)", bullets=[
    "Doença com prevalência 0,01%; sensibilidade 99%; especificidade 95%.",
    "O teste deu positivo. Qual a probabilidade de a pessoa ter a doença?",
], formula=r"P(D\mid +) = \frac{0{,}99\cdot0{,}0001}{0{,}99\cdot0{,}0001 + 0{,}05\cdot0{,}9999} \approx 0{,}2\%",
   example="Mesmo com um teste ótimo, a doença é tão rara que quase todos os positivos são falsos. Ainda assim, a chance subiu de 0,01% para 0,2% (20×).")
content("7 · Exercícios", "Megasena", bullets=[
    "Qual a probabilidade de ganhar com o jogo mínimo (6 números entre 60)?",
], formula=r"P = \frac{1}{\binom{60}{6}} = \frac{1}{50\,063\,860} \approx 2\times10^{-8}",
   example="Cerca de 0,000002%. Pode-se obter o mesmo encadeando 6/60 · 5/59 · ... · 1/55.")
content("7 · Exercícios", "Distribuição binomial", bullets=[
    "Prova com 10 questões, 5 alternativas cada, só uma correta.",
    "Qual a chance de passar (nota ≥ 5) apenas no chute?",
], formula=r"P(X\geq 5) = 1 - P(X\leq 4),\quad X\sim \mathrm{Bin}(10,\ 0{,}2)",
   example="Cerca de 3,3%. A nota mais provável no chute é 2. Com 20 questões, a chance cai para ~0,26%.")
content("7 · Exercícios", "Distribuição normal", bullets=[
    "Pressão arterial ~ N(100, 10). Que percentual está entre 80 e 120 mmHg?",
    "Acima de qual valor está apenas 1% da população?",
], formula=r"z = \frac{120-100}{10} = 2 \;\Rightarrow\; 95{,}4\%",
   example="Entre 80 e 120: ~95,4% (±2σ). O percentil 99 fica em 100 + 2,33·10 ≈ 123,3 mmHg.")
content("7 · Exercícios", "Intervalo de confiança", bullets=[
    "Amostra: 5, 3, 4, 2, 3, 4, 2, 3, 4, 5 (média 3,5; s ≈ 1,08).",
    "Construa o IC 95% para a média.",
], formula=r"3{,}5 \pm 2{,}262\cdot\frac{1{,}08}{\sqrt{10}} \approx [2{,}7;\ 4{,}3]",
   example="Se repetíssemos o estudo muitas vezes, ~95% dos intervalos conteriam a média da população.")
content("7 · Exercícios", "Teste de dois grupos", bullets=[
    "Besouros em duas florestas — Amostra 1: 8,12,15,21,25,44,44,60; Amostra 2: 2,4,5,9,12,17,19.",
    "Que teste usar?",
], example="Contagens baixas → distribuição assimétrica → Mann–Whitney (não paramétrico, bilateral, não pareado). Com p = 2,8% < 5%, a diferença é significativa.")
content("7 · Exercícios", "Qui-quadrado", bullets=[
    "Um dado lançado 60 vezes: frequências 8, 11, 7, 12, 15, 7.",
    "O dado é honesto?",
], formula=r"\chi^2 = \sum\frac{(O_i-10)^2}{10} = 5{,}2,\quad p \approx 0{,}39",
   example="Não se rejeita a hipótese de dado honesto. Mas, multiplicando tudo por 10, χ² = 52 e p ≈ 10⁻¹⁰: o teste depende do n absoluto.")
content("7 · Exercícios", "ANOVA", bullets=[
    "Contagem de linfócitos sob placebo e duas drogas, em 7 ninhadas.",
    "Animais da mesma ninhada; queremos saber se as drogas diferem do placebo.",
], example="Ninhada cria pareamento → ANOVA de medidas repetidas; pós-teste de Dunnett (contra controle). Resultado: p ≈ 0,002; só a droga A difere do placebo.")
content("7 · Exercícios", "Correlação e regressão", bullets=[
    "Altura × peso de 9 pessoas (152–193 cm; 38–89 kg).",
    "Calcule r, R² e a reta de regressão.",
], formula=r"r \approx 0{,}77,\quad R^2 \approx 0{,}59,\quad \hat{y} = 0{,}89x - 91",
   example="59% da variação do peso é explicada pela altura. A reta descreve a tendência, mas não prevê bem um indivíduo.")

refs_slide()

prs.save(str(OUT))
print(f"Slides salvos: {OUT.name}  ({len(prs.slides._sldIdLst)} slides)")
