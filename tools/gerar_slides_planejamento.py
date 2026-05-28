"""Deck de Planejamento Experimental — 20 slides (+ 5 sobre escolha do teste).
Baseado no conteúdo do curso IBCCF/UFRJ.
Gera: planejamento_experimental_slides.pptx"""

import hashlib
from pathlib import Path
from PIL import Image
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
from scipy import stats
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

BASE = Path(__file__).parent.parent
FIG  = BASE / "assets" / "slides"
DESN = FIG / "desenho"
FORM = FIG / "formulas"
FORM.mkdir(parents=True, exist_ok=True)
FUNDO = FIG / "fundo.png"
LOGO  = BASE / "assets" / "ibccf-logo.png"
OUT   = BASE / "planejamento_experimental_slides.pptx"

BG    = RGBColor(0xF3, 0xEC, 0xDF)
PAPER = RGBColor(0xFF, 0xFD, 0xF8)
INK   = RGBColor(0x1A, 0x1A, 0x1A)
BLUE  = RGBColor(0x32, 0x66, 0xAD)
RED   = RGBColor(0xC0, 0x39, 0x2B)
GREEN = RGBColor(0x1A, 0x7A, 0x4A)
MUTED = RGBColor(0x6B, 0x64, 0x57)
LINE  = RGBColor(0xC9, 0xBF, 0xA9)
TINT  = RGBColor(0xEC, 0xF1, 0xF8)
TINTG = RGBColor(0xDC, 0xEF, 0xE4)
TINTR = RGBColor(0xF6, 0xDE, 0xDB)

# paleta da tabela de escolha do teste ─────────────────────────────────
TC_NORM   = RGBColor(0xDC, 0xE7, 0xF4)   # azul-claro (paramétrico)
TC_NORM_K = RGBColor(0x32, 0x66, 0xAD)
TC_NONP   = RGBColor(0xFD, 0xEB, 0xD0)   # pêssego (não-paramétrico)
TC_NONP_K = RGBColor(0xC9, 0x6E, 0x1F)
TC_NOM    = RGBColor(0xDC, 0xEF, 0xE4)   # verde-claro (categóricas)
TC_NOM_K  = RGBColor(0x1A, 0x7A, 0x4A)
TC_SURV   = RGBColor(0xE8, 0xDA, 0xEF)   # roxo-claro (sobrevida)
TC_SURV_K = RGBColor(0x6B, 0x3F, 0xA0)
TC_PRED   = RGBColor(0xFC, 0xF3, 0xCF)   # amarelo-claro (regressões)
TC_PRED_K = RGBColor(0xA0, 0x7C, 0x1A)
TC_CORR   = RGBColor(0xF6, 0xDE, 0xDB)   # rosa-claro (correlação)
TC_CORR_K = RGBColor(0xC0, 0x39, 0x2B)

SERIF = "Georgia"
MONO  = "Consolas"

# multiplicador global de fontes nos slides
SCALE = 1.15

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK  = prs.slide_layouts[6]
_num   = 0


# ── fórmulas ──────────────────────────────────────────────────────────────
def formula_png(latex, fontsize=30, color="#1a1a1a"):
    key  = hashlib.md5(f"{latex}{fontsize}{color}".encode()).hexdigest()[:12]
    path = FORM / f"f_{key}.png"
    if not path.exists():
        fig = plt.figure(figsize=(0.01, 0.01))
        fig.text(0, 0, f"${latex}$", fontsize=fontsize, color=color)
        fig.savefig(path, dpi=200, transparent=True, bbox_inches="tight", pad_inches=0.06)
        plt.close(fig)
    return path


# ── helpers de layout ────────────────────────────────────────────────────
def _bg(slide):
    slide.shapes.add_picture(str(FUNDO), 0, 0, SW, SH)
    pic = slide.shapes[-1]
    slide.shapes._spTree.remove(pic._element)
    slide.shapes._spTree.insert(2, pic._element)


def _box(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    return tb


def _set(p, text, size, color=INK, bold=False, font=SERIF, italic=False):
    p.text = text
    for r in p.runs:
        r.font.size = Pt(int(round(size * SCALE)))
        r.font.color.rgb = color
        r.font.bold  = bold
        r.font.name  = font
        r.font.italic = italic


def _rect(slide, l, t, w, h, color, rounded=False):
    shp = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    sh  = slide.shapes.add_shape(shp, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background(); sh.shadow.inherit = False
    return sh


def _img_fit(slide, path, l, t, w, h):
    iw, ih  = Image.open(path).size
    ar, bar = iw / ih, w / h
    if ar > bar:
        nw = w; nh = int(w / ar)
    else:
        nh = h; nw = int(h * ar)
    slide.shapes.add_picture(str(path), l + (w - nw) // 2, t + (h - nh) // 2, nw, nh)
    return nw, nh


def _footer(slide):
    global _num
    _num += 1
    tb = _box(slide, Inches(0.5), Inches(7.04), Inches(6), Inches(0.4))
    _set(tb.text_frame.paragraphs[0],
         "Planejamento Experimental · IBCCF · UFRJ", 10, MUTED, font=MONO)
    nb = _box(slide, Inches(12.3), Inches(7.04), Inches(0.8), Inches(0.4))
    p  = nb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    _set(p, str(_num), 10, MUTED, font=MONO)


def _bullets(slide, items, l, t, w, h, size=18):
    tb = _box(slide, l, t, w, h)
    tf = tb.text_frame
    for i, it in enumerate(items):
        p   = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        sub = it.startswith("- ")
        txt = it[2:] if sub else it
        _set(p, ("–  " if sub else "•  ") + txt, size - (2 if sub else 0), INK)
        p.space_after  = Pt(9)
        p.line_spacing = 1.12
        if sub: p.level = 1


def _formula_card(slide, latex, cx_l, cx_w, top, target_h_in=0.7, fontsize=30):
    path = formula_png(latex, fontsize=fontsize)
    iw, ih = Image.open(path).size
    h = Inches(target_h_in); w = int(h * iw / ih)
    maxw = cx_w - Inches(0.8)
    if w > maxw:
        w = maxw; h = int(w * ih / iw)
    card_w = w + Inches(0.7); card_h = h + Inches(0.45)
    card_l = cx_l + (cx_w - card_w) // 2
    _rect(slide, card_l, top, card_w, card_h, PAPER, rounded=True)
    slide.shapes.add_picture(str(path),
                             card_l + (card_w - w) // 2,
                             top    + (card_h - h) // 2, w, h)
    return card_h


def _example(slide, text, l, t, w):
    _rect(slide, l, t, w, Inches(1.15), TINT, rounded=True)
    tb = _box(slide, l + Inches(0.2), t + Inches(0.08),
              w - Inches(0.4), Inches(1.0))
    tf = tb.text_frame
    _set(tf.paragraphs[0], "EXEMPLO", 11, BLUE, bold=True, font=MONO)
    p = tf.add_paragraph(); _set(p, text, 14, INK); p.line_spacing = 1.1


# ── tipos de slide ────────────────────────────────────────────────────────
def title_slide():
    s = prs.slides.add_slide(BLANK); _bg(s)
    if LOGO.exists():
        _img_fit(s, LOGO, Inches(4.67), Inches(1.0), Inches(4.0), Inches(1.25))
    tb = _box(s, Inches(1), Inches(2.7), Inches(11.33), Inches(1.6))
    p  = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    _set(p, "Planejamento Experimental", 54, INK, bold=True)
    sub = _box(s, Inches(1), Inches(4.15), Inches(11.33), Inches(0.6))
    p   = sub.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    _set(p, "Instituto de Biofísica Carlos Chagas Filho · UFRJ", 18, MUTED, font=MONO)
    au = _box(s, Inches(1), Inches(5.1), Inches(11.33), Inches(0.8))
    p  = au.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    _set(p, "Pedro Torres   ·   Gilberto Weissmuller", 18, INK)


def section_slide(num, title, subtitle=""):
    s = prs.slides.add_slide(BLANK); _bg(s)
    _rect(s, Inches(1), Inches(3.0), Inches(1.5), Inches(0.1), BLUE)
    nb = _box(s, Inches(1), Inches(2.0), Inches(4), Inches(1.0))
    _set(nb.text_frame.paragraphs[0], f"Parte {num}", 22, BLUE, font=MONO)
    tb = _box(s, Inches(1), Inches(3.25), Inches(11.3), Inches(1.8))
    _set(tb.text_frame.paragraphs[0], title, 44, INK, bold=True)
    if subtitle:
        sb = _box(s, Inches(1), Inches(5.1), Inches(11.3), Inches(0.8))
        _set(sb.text_frame.paragraphs[0], subtitle, 20, MUTED, italic=True)
    _footer(s)


def content(eyebrow, title, bullets=None, formula=None,
            image=None, caption=None, example=None):
    s  = prs.slides.add_slide(BLANK); _bg(s)
    eb = _box(s, Inches(0.7), Inches(0.42), Inches(11), Inches(0.4))
    _set(eb.text_frame.paragraphs[0], eyebrow.upper(), 12, MUTED, font=MONO)
    tb = _box(s, Inches(0.7), Inches(0.78), Inches(12), Inches(1.0))
    _set(tb.text_frame.paragraphs[0], title, 28, INK, bold=True)
    _rect(s, Inches(0.72), Inches(1.7), Inches(1.1), Inches(0.05), BLUE)
    top = Inches(2.05)

    if image and bullets:
        _bullets(s, bullets, Inches(0.55), Inches(1.95), Inches(5.6), Inches(5.0), size=20)
        _img_fit(s, image, Inches(6.2), Inches(1.85), Inches(7.0), Inches(5.0))
    elif image:
        _img_fit(s, image, Inches(0.30), Inches(1.85), Inches(12.73), Inches(5.0))
        if caption:
            cb = _box(s, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.3))
            p  = cb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            _set(p, caption, 13, MUTED, italic=True)
    elif bullets and formula:
        _bullets(s, bullets, Inches(0.9), top, Inches(11.5), Inches(3.0), size=19)
        y = top + Inches(3.0)
        _formula_card(s, formula, Inches(0.9), Inches(11.5), y, target_h_in=0.85)
    elif formula:
        _formula_card(s, formula, Inches(1.5), Inches(10.3), Inches(3.0),
                      target_h_in=1.20, fontsize=34)
    elif bullets:
        _bullets(s, bullets, Inches(0.9), top, Inches(11.5), Inches(4.6), size=20)

    if example:
        _example(s, example, Inches(0.9), Inches(5.85), Inches(11.5))
    _footer(s)


# ── slide de tabela 2 colunas ─────────────────────────────────────────────
def two_col_table(eyebrow, title, headers, rows, col1_w=4.5):
    """Tabela simples de 2 colunas."""
    s  = prs.slides.add_slide(BLANK); _bg(s)
    eb = _box(s, Inches(0.7), Inches(0.42), Inches(11), Inches(0.4))
    _set(eb.text_frame.paragraphs[0], eyebrow.upper(), 12, MUTED, font=MONO)
    tb = _box(s, Inches(0.7), Inches(0.78), Inches(12), Inches(1.0))
    _set(tb.text_frame.paragraphs[0], title, 28, INK, bold=True)
    _rect(s, Inches(0.72), Inches(1.7), Inches(1.1), Inches(0.05), BLUE)

    c1w   = Inches(col1_w)
    c2w   = Inches(13.333 - col1_w - 1.4)
    lm    = Inches(0.7)
    row_h = Inches(0.64)
    y     = Inches(1.95)

    # cabeçalho
    for ci, (hdr, cw) in enumerate(zip(headers, [c1w, c2w])):
        xl = lm + (c1w if ci else 0)
        _rect(s, xl, y, cw, row_h, BLUE)
        ht = _box(s, xl + Inches(0.12), y + Inches(0.1),
                  cw - Inches(0.2), row_h - Inches(0.14))
        _set(ht.text_frame.paragraphs[0], hdr, 17, PAPER, bold=True)
    y += row_h

    for ri, row in enumerate(rows):
        fc = RGBColor(0xEC, 0xE4, 0xD3) if ri % 2 == 0 else PAPER
        for ci, (cell, cw) in enumerate(zip(row, [c1w, c2w])):
            xl = lm + (c1w if ci else 0)
            _rect(s, xl, y, cw, row_h, fc)
            ct = _box(s, xl + Inches(0.12), y + Inches(0.09),
                      cw - Inches(0.2), row_h - Inches(0.1))
            _set(ct.text_frame.paragraphs[0], cell, 16, INK)
        y += row_h

    _footer(s)


# ── slide de quadro de 4 pilares (inline, sem figura externa) ─────────────
def four_boxes(eyebrow, title, items):
    """4 caixas coloridas horizontais com título e corpo."""
    s  = prs.slides.add_slide(BLANK); _bg(s)
    eb = _box(s, Inches(0.7), Inches(0.42), Inches(11), Inches(0.4))
    _set(eb.text_frame.paragraphs[0], eyebrow.upper(), 12, MUTED, font=MONO)
    tb = _box(s, Inches(0.7), Inches(0.78), Inches(12), Inches(1.0))
    _set(tb.text_frame.paragraphs[0], title, 28, INK, bold=True)
    _rect(s, Inches(0.72), Inches(1.7), Inches(1.1), Inches(0.05), BLUE)

    colors = [BLUE, RED, GREEN, MUTED]
    fills  = [
        RGBColor(0xDC, 0xE7, 0xF4),
        RGBColor(0xF6, 0xDE, 0xDB),
        RGBColor(0xDC, 0xEF, 0xE4),
        RGBColor(0xEC, 0xE4, 0xD3),
    ]
    bw = Inches(2.9); bh = Inches(3.8)
    gap = Inches(0.27)
    start_x = Inches(0.72)
    top     = Inches(2.05)

    for i, (titulo, corpo) in enumerate(items):
        xl = start_x + i * (bw + gap)
        sh = prs.slides[-1].shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, xl, top, bw, bh)
        sh.fill.solid(); sh.fill.fore_color.rgb = fills[i]
        sh.line.color.rgb = colors[i]; sh.line.width = Pt(2)
        sh.shadow.inherit = False

        ht = _box(prs.slides[-1], xl + Inches(0.2), top + Inches(0.25),
                  bw - Inches(0.4), Inches(0.7))
        _set(ht.text_frame.paragraphs[0], titulo, 17, colors[i], bold=True)

        ct = _box(prs.slides[-1], xl + Inches(0.2), top + Inches(1.1),
                  bw - Inches(0.4), bh - Inches(1.2))
        ct.text_frame.word_wrap = True
        _set(ct.text_frame.paragraphs[0], corpo, 14, INK)
        ct.text_frame.paragraphs[0].line_spacing = 1.2

    _footer(prs.slides[-1])


# ── tabela ÚNICA e completa de escolha de testes (dois painéis lado a lado) ──
def comprehensive_test_table():
    """Uma só tabela em dois painéis horizontais cobrindo todos os testes."""
    s  = prs.slides.add_slide(BLANK); _bg(s)
    eb = _box(s, Inches(0.65), Inches(0.30), Inches(12), Inches(0.35))
    _set(eb.text_frame.paragraphs[0], "COMO ESCOLHER O TESTE", 12, MUTED, font=MONO)
    tb = _box(s, Inches(0.65), Inches(0.60), Inches(12), Inches(0.80))
    _set(tb.text_frame.paragraphs[0],
         "Guia completo de escolha do teste estatístico", 26, INK, bold=True)
    _rect(s, Inches(0.67), Inches(1.40), Inches(1.1), Inches(0.05), BLUE)

    DATA_FS   = 15      # fonte nas células de dados
    HDR_FS    = 15      # fonte nos cabeçalhos de coluna
    SEC_FS    = 14      # fonte nos cabeçalhos de seção
    ROW_H     = Inches(0.385)   # altura de cada linha de dado
    SEC_H     = Inches(0.405)   # altura da linha de seção
    COL_HDR_H = Inches(0.42)    # altura do cabeçalho de coluna
    PAD       = Inches(0.11)    # padding horizontal interno

    # ── painel esquerdo: Diferenças ───────────────────────────────────────
    # colunas: Desfecho (3.0 in) | N (0.82 in) | Teste (2.88 in)  → total 6.70
    LX     = Inches(0.65)
    L_COLS = [Inches(3.0), Inches(0.82), Inches(2.88)]
    L_TOT  = sum(L_COLS)

    # ── painel direito: Correlação + Predição ─────────────────────────────
    # colunas: Desfecho / situação (2.80 in) | Teste (3.18 in)  → total 5.98
    GAP    = Inches(0.30)
    RX     = LX + L_TOT + GAP
    R_COLS = [Inches(2.80), Inches(3.18)]
    R_TOT  = sum(R_COLS)

    Y0 = Inches(1.53)   # topo dos painéis

    def _hdr_row(x, y, cols, texts, bg):
        cx = x
        for i, (t, cw) in enumerate(zip(texts, cols)):
            _rect(s, cx, y, cw, COL_HDR_H, bg)
            ht = _box(s, cx + PAD, y + Inches(0.06),
                      cw - 2*PAD, COL_HDR_H - Inches(0.08))
            ht.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            _set(ht.text_frame.paragraphs[0], t, HDR_FS, PAPER, bold=True)
            cx += cw
        return y + COL_HDR_H

    def _sec_row(x, y, total_w, text, bg, tc=PAPER):
        _rect(s, x, y, total_w, SEC_H, bg)
        ht = _box(s, x + PAD, y + Inches(0.06),
                  total_w - 2*PAD, SEC_H - Inches(0.08))
        _set(ht.text_frame.paragraphs[0], text, SEC_FS, tc, bold=True)
        return y + SEC_H

    def _data_row(x, y, cols, cells, fills, text_colors, bold_last=True):
        cx = x
        for i, (cell, cw, fc, tc) in enumerate(zip(cells, cols, fills, text_colors)):
            _rect(s, cx, y, cw, ROW_H, fc)
            ct = _box(s, cx + PAD, y + Inches(0.05),
                      cw - 2*PAD, ROW_H - Inches(0.06))
            ct.text_frame.paragraphs[0].alignment = (
                PP_ALIGN.CENTER if i == len(cells) - 1 else PP_ALIGN.LEFT)
            _set(ct.text_frame.paragraphs[0], cell, DATA_FS, tc,
                 bold=(bold_last and i == len(cells) - 1))
            cx += cw
        return y + ROW_H

    # ── PAINEL ESQUERDO ────────────────────────────────────────────────────
    y = _hdr_row(LX, Y0, L_COLS,
                 ["Tipo de desfecho", "N", "Teste"], INK)

    # seção independentes
    y = _sec_row(LX, y, L_TOT, "  Grupos independentes (não pareados)", TC_NORM_K)
    ind_rows = [
        ("Contínua — Normal",              "2",  "Teste t (Student / Welch)",
         TC_NORM, TC_NORM, TC_NORM_K,   INK, INK, PAPER),
        ("Contínua — Normal",              ">2", "ANOVA one-way",
         TC_NORM, TC_NORM, TC_NORM_K,   INK, INK, PAPER),
        ("Contínua não-Normal / Ordinal",  "2",  "Mann–Whitney U",
         TC_NONP, TC_NONP, TC_NONP_K,   INK, INK, PAPER),
        ("Contínua não-Normal / Ordinal",  ">2", "Kruskal–Wallis H",
         TC_NONP, TC_NONP, TC_NONP_K,   INK, INK, PAPER),
        ("Nominal (categorias)",           "2+", "Qui-quadrado · Fisher exato",
         TC_NOM,  TC_NOM,  TC_NOM_K,    INK, INK, PAPER),
        ("Sobrevida (tempo até evento)",   "—",  "Log-Rank · Kaplan–Meier",
         TC_SURV, TC_SURV, TC_SURV_K,   INK, INK, PAPER),
    ]
    for d, n, t, f0, f1, f2, c0, c1, c2 in ind_rows:
        y = _data_row(LX, y, L_COLS, [d, n, t], [f0, f1, f2], [c0, c1, c2])

    # seção pareados
    y = _sec_row(LX, y, L_TOT, "  Grupos pareados (mesmo indivíduo / blocos)", TC_CORR_K)
    par_rows = [
        ("Contínua — Normal",              "2",  "Teste t pareado",
         TC_NORM, TC_NORM, TC_NORM_K,   INK, INK, PAPER),
        ("Contínua — Normal",              ">2", "ANOVA medidas repetidas",
         TC_NORM, TC_NORM, TC_NORM_K,   INK, INK, PAPER),
        ("Contínua não-Normal / Ordinal",  "2",  "Wilcoxon dos postos",
         TC_NONP, TC_NONP, TC_NONP_K,   INK, INK, PAPER),
        ("Contínua não-Normal / Ordinal",  ">2", "Friedman",
         TC_NONP, TC_NONP, TC_NONP_K,   INK, INK, PAPER),
        ("Nominal (categorias)",           "2",  "McNemar",
         TC_NOM,  TC_NOM,  TC_NOM_K,    INK, INK, PAPER),
    ]
    for d, n, t, f0, f1, f2, c0, c1, c2 in par_rows:
        y = _data_row(LX, y, L_COLS, [d, n, t], [f0, f1, f2], [c0, c1, c2])

    # ── PAINEL DIREITO ─────────────────────────────────────────────────────
    y = _hdr_row(RX, Y0, R_COLS, ["Situação / desfecho", "Teste / modelo"], INK)

    # seção correlação
    y = _sec_row(RX, y, R_TOT, "  Correlação bivariada", TC_CORR_K)
    corr_rows = [
        ("Normal e linear",              "Pearson  (r)",
         TC_CORR, TC_CORR_K,  INK, PAPER),
        ("Não-Normal / Ordinal / monotônica", "Spearman (ρ)  ·  Kendall (τ)",
         TC_CORR, TC_CORR_K,  INK, PAPER),
    ]
    for d, t, f0, f1, c0, c1 in corr_rows:
        y = _data_row(RX, y, R_COLS, [d, t], [f0, f1], [c0, c1])

    # seção predição / regressão
    y = _sec_row(RX, y, R_TOT, "  Predição multivariável", TC_PRED_K,
                 tc=RGBColor(0xFF, 0xFF, 0xFF))
    pred_rows = [
        ("Desfecho contínuo",              "Regressão linear",
         TC_PRED, TC_PRED_K,  INK, PAPER),
        ("Desfecho ordinal",               "Regressão logística ordinal",
         TC_PRED, TC_PRED_K,  INK, PAPER),
        ("Desfecho nominal — 2 níveis",    "Regressão logística binária",
         TC_PRED, TC_PRED_K,  INK, PAPER),
        ("Desfecho nominal — >2 níveis",   "Regressão logística multinomial",
         TC_PRED, TC_PRED_K,  INK, PAPER),
        ("Sobrevida (tempo até evento)",   "Regressão de Cox",
         TC_SURV, TC_SURV_K,  INK, PAPER),
        ("Contagem / taxa",                "Regressão de Poisson",
         TC_PRED, TC_PRED_K,  INK, PAPER),
    ]
    for d, t, f0, f1, c0, c1 in pred_rows:
        y = _data_row(RX, y, R_COLS, [d, t], [f0, f1], [c0, c1])

    _footer(s)


# ── slide das 5 perguntas ─────────────────────────────────────────────────
def five_questions_slide(eyebrow, title):
    s  = prs.slides.add_slide(BLANK); _bg(s)
    eb = _box(s, Inches(0.7), Inches(0.42), Inches(11), Inches(0.4))
    _set(eb.text_frame.paragraphs[0], eyebrow.upper(), 12, MUTED, font=MONO)
    tb = _box(s, Inches(0.7), Inches(0.78), Inches(12), Inches(1.0))
    _set(tb.text_frame.paragraphs[0], title, 28, INK, bold=True)
    _rect(s, Inches(0.72), Inches(1.7), Inches(1.1), Inches(0.05), BLUE)

    perguntas = [
        ("Q1", "Bivariada ou multivariável?",
         "uma única exposição contra o desfecho, ou ajuste para várias covariáveis?"),
        ("Q2", "Diferença ou correlação?",
         "comparar grupos / medir um efeito, ou medir a associação entre variáveis?"),
        ("Q3", "Independente ou pareada?",
         "amostras separadas, ou o mesmo indivíduo (antes–depois, blocos)?"),
        ("Q4", "Tipo do desfecho — e normalidade?",
         "contínua (Normal ou não), ordinal, nominal, tempo até evento, contagem?"),
        ("Q5", "Quantos grupos ou condições?",
         "duas, ou três ou mais?"),
    ]
    cores = [TC_NORM_K, TC_CORR_K, TC_NOM_K, TC_SURV_K, TC_NONP_K]

    y0       = Inches(2.0)
    card_h   = Inches(0.85)
    card_gap = Inches(0.12)

    for i, ((q, head, body), col) in enumerate(zip(perguntas, cores)):
        y = y0 + i * (card_h + card_gap)
        # círculo com Qn
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                  Inches(0.7), y + Inches(0.1),
                                  Inches(0.7), Inches(0.7))
        circ.fill.solid(); circ.fill.fore_color.rgb = col
        circ.line.fill.background(); circ.shadow.inherit = False
        qt = _box(s, Inches(0.7), y + Inches(0.18), Inches(0.7), Inches(0.6))
        qt.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        _set(qt.text_frame.paragraphs[0], q, 19, PAPER, bold=True, font=MONO)
        # cartão
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(1.6), y, Inches(11.2), card_h)
        card.fill.solid(); card.fill.fore_color.rgb = PAPER
        card.line.color.rgb = col; card.line.width = Pt(1.6)
        card.shadow.inherit = False
        # cabeçalho da pergunta
        ht = _box(s, Inches(1.85), y + Inches(0.08),
                  Inches(10.9), Inches(0.45))
        _set(ht.text_frame.paragraphs[0], head, 20, col, bold=True)
        # corpo da pergunta
        bt = _box(s, Inches(1.85), y + Inches(0.50),
                  Inches(10.9), Inches(0.40))
        _set(bt.text_frame.paragraphs[0], body, 15, MUTED, italic=True)

    _footer(s)


# ── slide só com uma imagem grande ────────────────────────────────────────
def image_slide(eyebrow, title, image_path, caption=None):
    """Slide com imagem ocupando o máximo de espaço possível."""
    s  = prs.slides.add_slide(BLANK); _bg(s)
    eb = _box(s, Inches(0.7), Inches(0.30), Inches(11), Inches(0.4))
    _set(eb.text_frame.paragraphs[0], eyebrow.upper(), 12, MUTED, font=MONO)
    tb = _box(s, Inches(0.7), Inches(0.60), Inches(12), Inches(0.85))
    _set(tb.text_frame.paragraphs[0], title, 28, INK, bold=True)
    _rect(s, Inches(0.72), Inches(1.46), Inches(1.1), Inches(0.05), BLUE)
    top    = Inches(1.62)
    bot    = Inches(6.78) if caption else Inches(7.00)
    img_h  = bot - top
    _img_fit(s, image_path, Inches(0.30), top, Inches(12.73), img_h)
    if caption:
        cb = _box(s, Inches(0.5), Inches(6.82), Inches(12.3), Inches(0.3))
        p  = cb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        _set(p, caption, 13, MUTED, italic=True)
    _footer(s)


# ── refs ──────────────────────────────────────────────────────────────────
def refs_slide():
    s  = prs.slides.add_slide(BLANK); _bg(s)
    tb = _box(s, Inches(0.7), Inches(0.7), Inches(11), Inches(1))
    _set(tb.text_frame.paragraphs[0], "Referências e contato", 28, INK, bold=True)
    _rect(s, Inches(0.72), Inches(1.6), Inches(1.1), Inches(0.05), BLUE)
    if LOGO.exists():
        _img_fit(s, LOGO, Inches(0.9), Inches(2.3), Inches(4.5), Inches(1.5))
    cb = _box(s, Inches(0.9), Inches(4.1), Inches(11), Inches(2.6))
    tf = cb.text_frame
    _set(tf.paragraphs[0], "Material teórico completo e interativo:", 18, INK)
    p = tf.add_paragraph()
    _set(p, "monteirotorres.github.io/biostat", 18, BLUE, font=MONO)
    p = tf.add_paragraph(); _set(p, " ", 8)
    p = tf.add_paragraph()
    _set(p, "Notebooks em Python (pandas · seaborn · scipy.stats)", 15, MUTED)
    p = tf.add_paragraph()
    _set(p, "Contato: monteirotorres@biof.ufrj.br", 15, MUTED, font=MONO)
    _footer(s)


# =====================================================================
#                      CONTEÚDO — 20 SLIDES
# =====================================================================

# 1 — Capa
title_slide()

# ── PARTE 1: FUNDAMENTOS ──────────────────────────────────────────────────

# 2 — Por que planejar?
content(
    "planejamento experimental",
    "Nenhuma análise conserta um experimento mal desenhado",
    bullets=[
        "O resultado de um estudo depende fundamentalmente de como os dados foram coletados",
        "Uma análise sofisticada não recupera informação que jamais foi registrada",
        "Planejar bem = definir antes da coleta: o que medir, quantos indivíduos, "
        "como alocar os tratamentos e quais vieses controlar",
        "O cálculo amostral e o controle de confundidores têm que vir antes do primeiro tubo",
        "- Estudos subdimensionados geram 'negativos' que nada provam e 'positivos' que superestimam o efeito",
    ],
    example="Dois grupos de pesquisa investigam o mesmo fármaco. "
            "Um planeja com 80 % de poder e detecta o efeito. "
            "O outro coleta 'o que deu' e publica um falso negativo.",
)

# 3 — Os quatro pilares (com figura)
content(
    "planejamento experimental",
    "Os quatro pilares de um bom experimento",
    image=DESN / "pilares.png",
    caption="Controle · Aleatorização · Repetição · Pareamento — os pilares do desenho experimental clássico",
)

# 4 — Controle
content(
    "pilar 1",
    "Controle: isolar o que se quer estudar",
    bullets=[
        "Manter constantes todas as variáveis que não estão sendo manipuladas",
        "Grupo controle recebe placebo ou condição padrão nas mesmas condições do grupo experimental",
        "Mesmo observador, mesmo equipamento, mesmo momento do dia — quando possível",
        "- Variações não controladas inflam a variabilidade e reduzem o poder",
        "Exemplo clássico: se um grupo é medido de manhã e o outro à tarde, "
        "qualquer diferença pode ser ritmo circadiano, não o tratamento",
    ],
    example="Experimento de comportamento animal: todos os camundongos do mesmo sexo, "
            "mesma gaiola, mesma ração, testados na mesma janela de tempo.",
)

# 5 — Aleatorização
content(
    "pilar 2",
    "Aleatorização: o único antídoto para o confundidor",
    bullets=[
        "Sortear quem recebe cada tratamento distribui aleatoriamente os confundidores não conhecidos",
        "Sem aleatorização, grupos podem diferir em fatores que o pesquisador não mediu",
        "- Exemplo: se pacientes mais graves tendem a receber o tratamento A, qualquer diferença pode ser gravidade",
        "Aleatorização em blocos: garante equilíbrio por subgrupos (sexo, centro, turno)",
        "Duplo-cego: participante e avaliador desconhecem o grupo — controla viés de expectativa",
    ],
    example="Ensaio clínico: randomização por envelope selado ou gerador de números. "
            "Estudos observacionais não podem fazer isso — por isso não provam causalidade.",
)

# 6 — Repetição
content(
    "pilar 3",
    "Repetição: o acaso se dilui com o n",
    bullets=[
        "Qualquer medida individual carrega variabilidade; a média de n medidas tem desvio reduzido por √n",
        "Repetição técnica: mesma amostra medida mais de uma vez (controla erro de medida)",
        "Repetição biológica: amostras de indivíduos diferentes (controla variabilidade biológica)",
        "- Os dois tipos não são intercambiáveis: três medidas do mesmo camundongo ≠ três camundongos",
        "O n do cálculo amostral sempre se refere a unidades experimentais independentes",
    ],
    example="qPCR: triplicatas técnicas são controle de pipetagem. "
            "Para inferência biológica, precisamos de animais (ou pacientes) independentes.",
)

# 7 — Pareamento
content(
    "pilar 4",
    "Pareamento: cada indivíduo como seu próprio controle",
    image=DESN / "pareamento.png",
    caption="Sem pareamento, a variabilidade entre indivíduos domina o sinal. "
            "Pareando, cada diferença individual é removida.",
)

# 8 — Planejamento fatorial: conceito
content(
    "desenhos especiais",
    "Planejamento fatorial: vários fatores ao mesmo tempo",
    image=DESN / "desenho_fatorial.png",
    bullets=[
        "Em vez de variar UM fator de cada vez, varia-se TODOS simultaneamente",
        "Notação 2 × 2: dois fatores, dois níveis cada → 4 grupos",
        "Generaliza para 2 × 3, 3 × 3, 2 × 2 × 2, …",
        "Vantagens:",
        "- mais EFICIENTE: o mesmo n responde a várias perguntas",
        "- revela INTERAÇÕES entre fatores",
        "Analisado por ANOVA fatorial (two-way, three-way)",
    ],
)

# 9 — Efeitos principais e interação
content(
    "desenhos especiais",
    "Efeitos principais e interação",
    image=DESN / "interacao.png",
    caption="Linhas paralelas → sem interação · linhas não paralelas → "
            "o efeito de um fator depende do nível do outro",
)

# ── PARTE 2: TIPOS DE ESTUDO ───────────────────────────────────────────────

# 8 — Seção
section_slide(2, "Tipos de estudo e variáveis de confusão",
              "Quando podemos afirmar causalidade?")

# 9 — Experimento vs. observacional
two_col_table(
    "tipos de estudo",
    "Experimento verdadeiro vs. estudo observacional",
    headers=["Característica", "Experimento controlado      |      Estudo observacional"],
    rows=[
        ("Variável manipulada?", "Sim, pelo pesquisador              |      Não — apenas observada"),
        ("Aleatorização?",       "Sim (RCT)                                 |      Não"),
        ("Permite causalidade?", "Sim                                           |      Não — apenas associação"),
        ("Controla confundidores?", "Por aleatorização              |      Parcialmente (propensity score, etc.)"),
        ("Exemplos",             "Ensaio clínico, RCBD                |      Coorte, caso-controle, transversal"),
    ],
    col1_w=4.0,
)

# 10 — Variáveis de confusão
content(
    "tipos de estudo",
    "Variável de confusão: o culpado escondido",
    bullets=[
        "Um confundidor está associado tanto à exposição quanto ao desfecho — e não é intermediário",
        "Sem controle, seu efeito se mistura ao do tratamento e distorce a estimativa",
        "Exemplo histórico: motoristas de ônibus tinham mais doença coronariana do que cobradores",
        "- A diferença real era sedentarismo (motoristas ficam sentados), não o emprego em si",
        "- A solução foi estratificar por tipo de atividade — e o efeito do exercício emergiu",
        "Aleatorização distribui confundidores conhecidos E desconhecidos — a vantagem do experimento",
    ],
    example="Café e câncer de pulmão: fumantes bebem mais café. "
            "O tabagismo é o confundidor; sem controlá-lo, o café parece cancerígeno.",
)

# ── PARTE 3: PODER E TAMANHO AMOSTRAL ─────────────────────────────────────

# 11 — Seção
section_slide(3, "Poder estatístico e cálculo amostral",
              "Quantos sujeitos preciso para responder minha pergunta?")

# 12 — Erros tipo I e tipo II
two_col_table(
    "poder estatístico",
    "Os dois tipos de erro e o que os controla",
    headers=["", "H₀ é verdadeira          |          H₀ é falsa"],
    rows=[
        ("Não rejeitamos H₀",
         "Decisão correta (1 − α)   |   Erro tipo II  (β)  — falso negativo"),
        ("Rejeitamos H₀",
         "Erro tipo I  (α)               |   Decisão correta  (poder = 1 − β)"),
        ("O que controla?",
         "Nível de significância α  |   n, efeito real, σ, α"),
        ("Consequência",
         "Alarme falso                        |   Perde efeito real — estudo inconclusivo"),
    ],
    col1_w=3.5,
)

# 13 — Poder estatístico
content(
    "poder estatístico",
    "Poder: a chance de detectar o que existe",
    bullets=[
        "Poder = 1 − β = P(rejeitar H₀ | H₁ é verdadeira)",
        "Convenção: queremos poder ≥ 80 % (β ≤ 20 %)",
        "Quatro fatores controlam o poder — fixados três, o quarto está determinado:",
        "- ↑ n  →  ↑ poder (mais informação = menos incerteza)",
        "- ↑ δ (efeito real)  →  ↑ poder (sinal maior é mais fácil de ver)",
        "- ↑ σ (variabilidade)  →  ↓ poder (ruído esconde o sinal)",
        "- ↑ α  →  ↑ poder, mas ↑ erro tipo I (tradeoff)",
    ],
    formula=r"\text{poder} = 1 - \beta = P\!\left(\text{rejeitar }H_0 \mid H_1\right)",
)

# 14 — Os quatro quadrantes
content(
    "poder estatístico",
    "Os quatro quadrantes: efeito × variabilidade",
    image=DESN / "quadrantes.png",
    caption="A combinação de tamanho de efeito e variabilidade determina o esforço amostral necessário",
)

# 15 — Cálculo amostral
content(
    "cálculo amostral",
    "Quantos sujeitos preciso?",
    bullets=[
        "A pergunta central do planejamento: quantos indivíduos para 80 % de poder, α = 0,05?",
        "Para comparação de duas médias (teste t, Welch):",
        "- d = δ/σ é o tamanho de efeito padronizado (Cohen's d)",
        "- z₀.₉₇₅ ≈ 1,96  e  z₀.₈₀ ≈ 0,84 para poder = 80 %",
        "Exemplo: d = 0,5 (efeito médio) → n ≈ 64 por grupo",
        "Exemplo: d = 0,2 (efeito pequeno) → n ≈ 394 por grupo",
        "- Dobrar a precisão exige quatro vezes mais sujeitos",
    ],
    formula=r"n \approx 2\!\left(\frac{z_{1-\alpha/2}+z_{1-\beta}}{d}\right)^{\!2} \;\text{ por grupo}",
)

# 16 — Curvas de poder
content(
    "cálculo amostral",
    "Poder em função do n: curvas para diferentes efeitos",
    image=DESN / "poder_curvas.png",
    caption="Efeito grande (d = 0,8): poucos sujeitos bastam. "
            "Efeito pequeno (d = 0,2): centenas. "
            "A linha tracejada marca poder = 80 %.",
)

# ── PARTE 4: SIGNIFICÂNCIA E RELEVÂNCIA ───────────────────────────────────

# 17 — Seção
section_slide(4, "Significância, relevância e tamanho de efeito",
              "Um p-valor pequeno não significa que o efeito importa")

# 18 — Significância vs. relevância
content(
    "significância vs. relevância",
    "Significativo não é o mesmo que importante",
    bullets=[
        "Significância estatística: a diferença é grande o bastante para ser detectada pelo design",
        "Relevância clínica/biológica: a diferença é grande o bastante para importar na prática",
        "Com n enorme qualquer diferença minúscula pode ser estatisticamente significativa",
        "Exemplo: diferença de 0,1 ponto de QI entre primogênitos e caçulas — p < 0,001 em estudo com 250 000 pares, mas clinicamente irrelevante",
        "- O que importa: o tamanho do efeito e seu intervalo de confiança",
        "Reporte sempre: valor-p + tamanho de efeito + IC 95 %",
    ],
    example="Antitérmico reduz febre em 0,1 °C: p = 0,02 com n grande. "
            "Estatisticamente significativo, clinicamente inútil.",
)

# 19 — Tamanhos de efeito padronizados
two_col_table(
    "tamanho de efeito",
    "Medidas de tamanho de efeito padronizadas",
    headers=["Medida", "Quando usar                          |   Pequeno   Médio   Grande"],
    rows=[
        ("Cohen's d",   "Diferença de duas médias (t)              |     0,2        0,5        0,8"),
        ("Cohen's f",   "ANOVA (variação entre grupos)          |     0,10      0,25      0,40"),
        ("Cohen's w",   "Qui-quadrado (tabelas de conting.)   |     0,10      0,30      0,50"),
        ("r de Pearson", "Correlação linear                                |     0,10      0,30      0,50"),
        ("η² (eta²)",   "Variância explicada pela ANOVA       |     0,01      0,06      0,14"),
    ],
    col1_w=3.2,
)

# ── PARTE 5: COMO ESCOLHER O TESTE ────────────────────────────────────────

# 20 — As cinco perguntas
five_questions_slide(
    "como escolher o teste",
    "Cinco perguntas que decidem o teste estatístico",
)

# 21 — tabela completa: imagem original do guia (em inglês)
image_slide(
    "como escolher o teste",
    "Guia completo de escolha do teste estatístico",
    DESN / "guia_testes.jpg",
    caption="A guide for choosing the most common statistical tests",
)

# 24 — Pós-testes e boas práticas
content(
    "depois do teste",
    "Pós-testes e o que sempre reportar",
    bullets=[
        "Uma ANOVA significativa diz que ALGUMA média difere — mas não qual",
        "- Tukey HSD: todos os grupos comparados entre si",
        "- Dunnett: cada grupo comparado com um único controle",
        "- Bonferroni · Šidák: poucas comparações planejadas a priori",
        "- Scheffé: combinações lineares de médias (contrastes)",
        "Não-paramétricos: Dunn (após Kruskal–Wallis) · Wilcoxon pareado corrigido (após Friedman)",
        "Antes de qualquer teste: faça um gráfico (histograma, boxplot, dispersão) e cheque suposições",
        "Sempre reporte valor-p + tamanho de efeito + intervalo de confiança",
    ],
)

# 25 — Checklist e referências
content(
    "checklist",
    "Checklist do bom experimento",
    bullets=[
        "Pergunta de pesquisa clara e hipótese pré-registrada antes da coleta",
        "Cálculo amostral feito com tamanho de efeito mínimo relevante",
        "Aleatorização documentada (como foi feita, por quem, quando)",
        "Grupos equilibrados em covariáveis importantes (baseline table)",
        "Cegamento do avaliador quando possível",
        "Análise dos dados cega ou pré-especificada no protocolo",
        "Reporte: valor-p, tamanho de efeito (+ IC 95 %) e poder pós-hoc se negativo",
        "- Um estudo negativo bem dimensionado diz tanto quanto um positivo",
    ],
)

# final: referências
refs_slide()

# ── salva ──────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f"OK  {OUT.relative_to(BASE)}  ({_num} slides)")
