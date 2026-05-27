"""Apresentação focada: Desenho Experimental.
Mesmo estilo dos demais slides (marca-d'água UFRJ·IBCCF, fórmulas em
cartões, paleta bege/azul)."""

import hashlib
from pathlib import Path
from PIL import Image
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

BASE = Path(__file__).parent.parent
FIG = BASE / "assets" / "slides"
FIG_D = FIG / "desenho"
FORM = FIG / "formulas"
FORM.mkdir(parents=True, exist_ok=True)
FUNDO = FIG / "fundo.png"
LOGO = BASE / "assets" / "ibccf-logo.png"
FLUXO = BASE / "assets" / "fluxograma_testes.png"
OUT = BASE / "desenho_experimental_slides.pptx"

PAPER = RGBColor(0xFF, 0xFD, 0xF8)
INK = RGBColor(0x1A, 0x1A, 0x1A)
BLUE = RGBColor(0x32, 0x66, 0xAD)
RED = RGBColor(0xC0, 0x39, 0x2B)
GREEN = RGBColor(0x1A, 0x7A, 0x4A)
MUTED = RGBColor(0x6B, 0x64, 0x57)
TINT = RGBColor(0xEC, 0xF1, 0xF8)

SERIF = "Georgia"; MONO = "Consolas"
prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]
_num = 0


def formula_png(latex, fontsize=30, color="#1a1a1a"):
    key = hashlib.md5(f"{latex}{fontsize}{color}".encode()).hexdigest()[:12]
    path = FORM / f"f_{key}.png"
    if not path.exists():
        fig = plt.figure(figsize=(0.01, 0.01))
        fig.text(0, 0, f"${latex}$", fontsize=fontsize, color=color)
        fig.savefig(path, dpi=200, transparent=True, bbox_inches="tight", pad_inches=0.06)
        plt.close(fig)
    return path


def _bg(slide):
    slide.shapes.add_picture(str(FUNDO), 0, 0, SW, SH)
    pic = slide.shapes[-1]
    slide.shapes._spTree.remove(pic._element)
    slide.shapes._spTree.insert(2, pic._element)


def _box(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(l, t, w, h); tb.text_frame.word_wrap = True; return tb


def _set(p, text, size, color=INK, bold=False, font=SERIF, italic=False):
    p.text = text
    for r in p.runs:
        r.font.size = Pt(size); r.font.color.rgb = color
        r.font.bold = bold; r.font.name = font; r.font.italic = italic


def _rect(slide, l, t, w, h, color, rounded=False):
    shp = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    sh = slide.shapes.add_shape(shp, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background(); sh.shadow.inherit = False
    return sh


def _img_fit(slide, path, l, t, w, h):
    iw, ih = Image.open(path).size
    ar = iw / ih; box_ar = w / h
    if ar > box_ar:
        nw = w; nh = int(w / ar)
    else:
        nh = h; nw = int(h * ar)
    slide.shapes.add_picture(str(path), l + (w - nw) // 2, t + (h - nh) // 2, nw, nh)


def _footer(slide):
    global _num; _num += 1
    tb = _box(slide, Inches(0.5), Inches(7.04), Inches(7), Inches(0.4))
    _set(tb.text_frame.paragraphs[0], "Desenho experimental · Bioestatística · IBCCF · UFRJ", 10, MUTED, font=MONO)
    nb = _box(slide, Inches(12.3), Inches(7.04), Inches(0.8), Inches(0.4))
    p = nb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    _set(p, str(_num), 10, MUTED, font=MONO)


def _bullets(slide, items, l, t, w, h, size=18):
    tb = _box(slide, l, t, w, h); tf = tb.text_frame
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        sub = it.startswith("- "); txt = it[2:] if sub else it
        _set(p, ("–  " if sub else "•  ") + txt, size - (2 if sub else 0), INK)
        p.space_after = Pt(9); p.line_spacing = 1.12
        if sub: p.level = 1


def _formula_card(slide, latex, cx_l, cx_w, top, target_h_in=0.7, fontsize=30):
    path = formula_png(latex, fontsize=fontsize)
    iw, ih = Image.open(path).size
    h = Inches(target_h_in); w = int(h * iw / ih)
    maxw = cx_w - Inches(0.8)
    if w > maxw: w = maxw; h = int(w * ih / iw)
    card_w = w + Inches(0.7); card_h = h + Inches(0.45)
    card_l = cx_l + (cx_w - card_w) // 2
    _rect(slide, card_l, top, card_w, card_h, PAPER, rounded=True)
    slide.shapes.add_picture(str(path), card_l + (card_w - w) // 2, top + (card_h - h) // 2, w, h)


def _example(slide, text, l, t, w, label="EXEMPLO", color=BLUE):
    _rect(slide, l, t, w, Inches(1.15), TINT, rounded=True)
    tb = _box(slide, l + Inches(0.2), t + Inches(0.08), w - Inches(0.4), Inches(1.0))
    tf = tb.text_frame
    _set(tf.paragraphs[0], label, 11, color, bold=True, font=MONO)
    p = tf.add_paragraph(); _set(p, text, 14, INK); p.line_spacing = 1.1


# ---------- tipos de slide ----------
def title_slide():
    s = prs.slides.add_slide(BLANK); _bg(s)
    if LOGO.exists():
        _img_fit(s, LOGO, Inches(4.67), Inches(1.0), Inches(4.0), Inches(1.25))
    tb = _box(s, Inches(0.5), Inches(2.7), Inches(12.3), Inches(1.4))
    _set(tb.text_frame.paragraphs[0], "Desenho Experimental", 56, INK, bold=True)
    tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    sub = _box(s, Inches(1), Inches(4.2), Inches(11.33), Inches(0.6))
    p = sub.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    _set(p, "Planejando experimentos antes de coletar dados", 20, MUTED, italic=True)
    sub2 = _box(s, Inches(1), Inches(4.9), Inches(11.33), Inches(0.5))
    p = sub2.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    _set(p, "Bioestatística · Instituto de Biofísica Carlos Chagas Filho · UFRJ", 14, MUTED, font=MONO)
    au = _box(s, Inches(1), Inches(5.7), Inches(11.33), Inches(0.6))
    p = au.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    _set(p, "Pedro Torres   ·   Gilberto Weissmuller", 16, INK)


def content(eyebrow, title, bullets=None, formula=None, image=None, image_dir=FIG_D,
            caption=None, example=None, ex_label="EXEMPLO", ex_color=BLUE):
    s = prs.slides.add_slide(BLANK); _bg(s)
    eb = _box(s, Inches(0.7), Inches(0.42), Inches(11), Inches(0.4))
    _set(eb.text_frame.paragraphs[0], eyebrow.upper(), 12, MUTED, font=MONO)
    tb = _box(s, Inches(0.7), Inches(0.78), Inches(12), Inches(1.0))
    _set(tb.text_frame.paragraphs[0], title, 28, INK, bold=True)
    _rect(s, Inches(0.72), Inches(1.7), Inches(1.1), Inches(0.05), BLUE)
    top = Inches(2.05)
    if image and bullets:
        _bullets(s, bullets, Inches(0.7), top, Inches(5.5), Inches(4.6))
        _img_fit(s, image_dir / image, Inches(6.5), top, Inches(6.3), Inches(4.4))
    elif image:
        _img_fit(s, image_dir / image, Inches(1.3), top, Inches(10.7), Inches(4.5))
        if caption:
            cb = _box(s, Inches(1), Inches(6.55), Inches(11.33), Inches(0.4))
            p = cb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            _set(p, caption, 12, MUTED, italic=True)
    elif bullets and formula:
        _bullets(s, bullets, Inches(0.9), top, Inches(11.5), Inches(3.0), size=19)
        _formula_card(s, formula, Inches(0.9), Inches(11.5), top + Inches(3.0), target_h_in=0.75)
    elif formula:
        _formula_card(s, formula, Inches(1.5), Inches(10.3), Inches(3.0), target_h_in=1.0, fontsize=34)
    elif bullets:
        _bullets(s, bullets, Inches(0.9), top, Inches(11.5), Inches(4.6), size=20)
    if example:
        _example(s, example, Inches(0.9), Inches(5.85), Inches(11.5), label=ex_label, color=ex_color)
    _footer(s)


def section_slide(num, title):
    s = prs.slides.add_slide(BLANK); _bg(s)
    _rect(s, Inches(1), Inches(3.05), Inches(1.5), Inches(0.12), BLUE)
    nb = _box(s, Inches(1), Inches(2.05), Inches(4), Inches(1.0))
    _set(nb.text_frame.paragraphs[0], f"Parte {num}", 22, BLUE, font=MONO)
    tb = _box(s, Inches(1), Inches(3.35), Inches(11.3), Inches(1.6))
    _set(tb.text_frame.paragraphs[0], title, 44, INK, bold=True)
    _footer(s)


def refs_slide():
    s = prs.slides.add_slide(BLANK); _bg(s)
    tb = _box(s, Inches(0.7), Inches(0.7), Inches(11), Inches(1))
    _set(tb.text_frame.paragraphs[0], "Material complementar", 28, INK, bold=True)
    _rect(s, Inches(0.72), Inches(1.6), Inches(1.1), Inches(0.05), BLUE)
    if LOGO.exists():
        _img_fit(s, LOGO, Inches(0.9), Inches(2.3), Inches(4.5), Inches(1.5))
    cb = _box(s, Inches(0.9), Inches(4.1), Inches(11), Inches(2.4))
    tf = cb.text_frame
    _set(tf.paragraphs[0], "Site do curso (versão interativa, com calculadora de potência):", 17, INK)
    p = tf.add_paragraph(); _set(p, "monteirotorres.github.io/biostat", 18, BLUE, font=MONO)
    p = tf.add_paragraph(); _set(p, " ", 8)
    p = tf.add_paragraph(); _set(p, "Notebooks em Python (pandas · seaborn · scipy.stats) acompanham cada tópico.", 14, MUTED)
    p = tf.add_paragraph(); _set(p, "Contato: monteirotorres@biof.ufrj.br", 14, MUTED, font=MONO)
    _footer(s)


# =====================================================================
#                              CONTEÚDO
# =====================================================================
title_slide()

# ---- 1. INTRODUÇÃO ----
section_slide(1, "Por que pensar antes de medir")

content("Desenho experimental", "O que é desenho experimental", bullets=[
    "É o planejamento das medidas: quem entra no estudo, em que condições e em qual sequência.",
    "Acontece antes da coleta — nenhuma análise estatística conserta um experimento mal desenhado.",
    "Determina o que se pode (e o que não se pode) concluir a partir dos dados.",
    "Garante que diferenças observadas reflitam o efeito de interesse, e não o ruído da coleta.",
])
content("Desenho experimental", "Por que vem antes da estatística", bullets=[
    "A análise correta de um experimento mal planejado ainda assim é inconclusiva.",
    "Um experimento subdimensionado expõe sujeitos sem chance real de chegar a uma resposta.",
    "Resultados positivos em estudos pequenos costumam superestimar o efeito real (winner's curse).",
    "O cálculo amostral e a escolha do teste devem ser feitos no projeto, não depois.",
])
content("Desenho experimental", "Os quatro pilares", image="pilares.png",
        caption="Controle, aleatorização, repetição e pareamento. Cada experimento usa, em maior ou menor grau, todos eles.")

# ---- 2. PILARES ----
section_slide(2, "Controle, aleatorização, repetição e pareamento")

content("Pilar 1 · Controle", "Manter as variáveis estranhas constantes", bullets=[
    "Tudo o que não está sendo estudado precisa, idealmente, ser igual entre os grupos.",
    "Equipamentos, reagentes, tempo de análise, condições ambientais.",
    "Quanto mais variáveis ficam soltas, mais difícil isolar o efeito de interesse.",
], example="Comparar a glicemia de pacientes em horários diferentes do dia, com balanças diferentes, em laboratórios diferentes — qualquer diferença observada pode vir do tratamento ou de qualquer um desses fatores.")
content("Pilar 2 · Aleatorização", "Sortear quem recebe cada tratamento", bullets=[
    "Distribui de forma equilibrada — em média — tanto as variáveis conhecidas quanto as desconhecidas.",
    "É a principal arma contra variáveis de confusão.",
    "Difere fundamentalmente do uso de voluntários, que se autosselecionam.",
    "Sem aleatorização, o estudo só revela associação; com ela, é possível inferir causa.",
])
content("Pilar 3 · Repetição", "Reduzir o efeito do acaso", bullets=[
    "Medir vários indivíduos no mesmo experimento — não basta uma única medida.",
    "A precisão da média melhora com a raiz do tamanho da amostra.",
    "Quadruplicar n reduz o erro padrão pela metade.",
    "Repetir um experimento, em outro laboratório, é o que valida descobertas científicas.",
], formula=r"\mathrm{EP} = \frac{\sigma}{\sqrt{n}}")
content("Pilar 4 · Pareamento", "Aproveitar variabilidade conhecida", image="pareamento.png",
        caption="Pareando, a variabilidade entre indivíduos sai da conta — o teste fica mais potente.")

# ---- 3. TIPOS DE ESTUDO ----
section_slide(3, "Tipos de estudo e variáveis de confusão")

content("Tipos de estudo", "Experimento verdadeiro vs. observacional", bullets=[
    "Experimento verdadeiro: o pesquisador manipula a variável independente e aleatoriza os tratamentos.",
    "Estudo observacional: os grupos já existem; o pesquisador apenas observa.",
    "Só o experimento aleatorizado permite afirmar uma relação de causa e efeito.",
    "O observacional revela associação — base para hipóteses, não para conclusões causais.",
])
content("Confundidores", "Variáveis de confusão", bullets=[
    "Uma variável é de confusão quando afeta o desfecho e, ao mesmo tempo, difere entre os grupos comparados.",
    "Sem aleatorização, qualquer característica desbalanceada pode explicar a diferença observada.",
    "Idade, sexo, gravidade da doença, estresse e saúde prévia são confundidores frequentes.",
    "A aleatorização os neutraliza, em média; o controle estatístico (por exemplo, ANCOVA) ajusta a posteriori.",
])
content("Confundidores", "Caso: motoristas e cobradores de ônibus", bullets=[
    "Estudo de Londres (1961): cobradores, que andavam o dia todo, tinham menos doença cardíaca que motoristas.",
    "É tentador concluir que o exercício é a causa — mas o estudo é observacional.",
    "Idade, saúde prévia e estresse do trabalho diferem entre os dois grupos.",
    "Anos depois descobriu-se que os motoristas já eram, em média, mais pesados ao serem contratados.",
], ex_label="LIÇÃO", ex_color=RED, example="Há associação real entre exercício e menos doença cardíaca — mas o estudo, sozinho, não prova causa.")

# ---- 4. AMOSTRAGEM ----
section_slide(4, "Técnicas de amostragem")

content("Amostragem", "Probabilística × não probabilística", bullets=[
    "Probabilística: cada indivíduo tem chance conhecida e não nula de ser escolhido.",
    "Permite estimar parâmetros populacionais e construir intervalos de confiança.",
    "Não probabilística (conveniência, voluntários, julgamento) gera viés de seleção.",
    "Útil em estudos exploratórios; nunca para comparar intervenções.",
])
content("Amostragem", "Aleatória simples e estratificada", bullets=[
    "Aleatória simples: todos têm a mesma chance — como num sorteio.",
    "Estratificada: a população é dividida em estratos homogêneos e sorteia-se dentro de cada um.",
    "Estratificar garante representação de subgrupos pequenos e reduz a variabilidade.",
    "Usada quando há subgrupos com características muito diferentes (faixas etárias, sexo, região).",
])
content("Amostragem", "Conglomerados e sistemática", bullets=[
    "Conglomerados: sorteiam-se grupos inteiros (escolas, bairros) e mede-se todos dentro deles.",
    "Barata logisticamente, mas menos precisa que a estratificada.",
    "Sistemática: escolhe-se um a cada k indivíduos de uma lista — prática, mas cuidado com padrões periódicos.",
])
content("Amostragem", "Cuidado com voluntários", bullets=[
    "Quem se voluntaria já tem, em média, características diferentes do grupo geral.",
    "Em estudos comparativos, isso vira viés de seleção — o resultado fica artificialmente positivo.",
    "A solução é sortear quem recebe cada tratamento, independentemente do interesse do participante.",
], ex_label="EXEMPLO", example="Hospital quer testar treinamento de lavagem das mãos. Se os 100 voluntários recebem o treinamento e são comparados com os demais, o estudo é inválido — os voluntários já são, em média, mais cuidadosos.")
content("Amostragem", "Pareamento em blocos", bullets=[
    "Quando há uma fonte conhecida de variabilidade, agrupar unidades semelhantes aumenta a potência.",
    "Cada bloco contém uma unidade de cada tratamento — distribuídas por sorteio dentro do bloco.",
    "A variação entre blocos não se confunde com o efeito do tratamento.",
], ex_label="EXEMPLO", example="Cinco ninhadas de ratos com filhotes parecidos entre si. Para testar uma droga em quatro doses, pegue um filhote de cada ninhada para cada dose — não misture os filhotes ao acaso.")

# ---- 5. ERROS E POTÊNCIA ----
section_slide(5, "Erros, potência e cálculo amostral")

content("Inferência", "Hipóteses e tipos de erro", bullets=[
    "H₀: ausência de efeito. H₁: existe efeito.",
    "Erro tipo I (α): rejeitar H₀ quando ela é verdadeira — falso positivo.",
    "Erro tipo II (β): não rejeitar H₀ quando ela é falsa — falso negativo.",
    "Fixamos α antes do estudo (geralmente 5%); β depende do n, do efeito e da variabilidade.",
])
content("Inferência", "Visualizando α e β", image="erros.png", image_dir=FIG,
        caption="α (vermelho) é o falso positivo sob H₀; β (azul) é o falso negativo sob H₁; o poder é 1 − β.")
content("Potência", "Poder estatístico", bullets=[
    "Probabilidade de detectar um efeito que realmente existe.",
    "Convenção: poder ≥ 80% — abaixo disso o estudo é subdimensionado.",
    "Quatro fatores estão amarrados: n, tamanho de efeito, variabilidade e α.",
    "Fixados três, o quarto está determinado.",
], formula=r"\mathrm{poder} = 1 - \beta")
content("Potência", "Curvas de poder × tamanho da amostra", image="poder_curvas.png",
        caption="Efeitos grandes precisam de pouca amostra; efeitos pequenos exigem n grande para chegar a 80%.")
content("Potência", "Os quatro quadrantes", image="quadrantes.png",
        caption="Combinação de tamanho de efeito e variabilidade dita a viabilidade do experimento.")
content("Potência", "Cálculo amostral", bullets=[
    "Quantos indivíduos preciso para detectar um efeito de tamanho δ com poder de 80% e α = 0,05?",
    "Depende do efeito esperado, da variabilidade e do nível de significância.",
    "Deve ser feito no planejamento e reportado no protocolo.",
    "Não há um número universal — é sempre uma estimativa para o problema específico.",
], formula=r"n \approx 2\left(\frac{z_{1-\alpha/2}+z_{1-\beta}}{\delta/\sigma}\right)^2")
content("Potência", "Calculadora interativa", image="calc_preview.png",
        caption="A versão online do curso traz uma calculadora para teste t, ANOVA, qui-quadrado e proporções.")
content("Potência", "Significância × relevância", bullets=[
    "Significância estatística: o efeito é detectável dado o desenho.",
    "Relevância prática: o efeito é grande o suficiente para importar.",
    "Com n enorme, é possível detectar efeitos minúsculos e irrelevantes.",
    "Reporte sempre tamanho de efeito e intervalo de confiança, não apenas o valor-p.",
], ex_label="LIÇÃO", ex_color=RED, example="Um efeito antitérmico de 0,1 °C pode ser estatisticamente detectável com n suficiente — e clinicamente inútil.")

# ---- 6. ESCOLHA DO TESTE ----
section_slide(6, "Escolhendo o teste e o pós-teste")

content("Análise", "Que teste escolher", image=Path("..") / "fluxograma_testes.png", image_dir=FIG,
        caption="Três perguntas: que tipo de dado, quantos grupos, e a normalidade vale?")
# nota: a imagem do fluxograma está em assets/, não em assets/slides/desenho/
content("Análise", "Equivalências paramétrico × não-paramétrico", bullets=[
    "1 amostra: teste t · Wilcoxon do sinal.",
    "2 independentes: teste t de Welch · Mann–Whitney.",
    "2 pareados: teste t pareado · Wilcoxon pareado.",
    "3+ independentes: ANOVA · Kruskal–Wallis.",
    "3+ pareados: ANOVA de medidas repetidas · Friedman.",
])
content("Análise", "Pós-testes da ANOVA", bullets=[
    "Todos os pares: Tukey HSD.",
    "Cada grupo contra um controle: Dunnett.",
    "Poucas comparações planejadas: Bonferroni ou Šidák.",
    "Combinações lineares de médias (contrastes): Scheffé.",
    "Após Kruskal–Wallis: teste de Dunn com correção de Bonferroni.",
])

# ---- 7. CHECKLIST ----
section_slide(7, "Checklist antes de começar")

content("Checklist", "Antes de coletar dados", bullets=[
    "Pergunta de pesquisa formulada de forma operacional.",
    "Variáveis dependentes e independentes claramente definidas, com unidades.",
    "Tamanho de efeito mínimo de interesse clínico/biológico definido.",
    "Cálculo amostral feito e justificado (idealmente com poder ≥ 80%).",
    "Teste estatístico escolhido antes da coleta — não após olhar os dados.",
    "Estratégia de aleatorização e cegamento descrita no protocolo.",
])
content("Checklist", "Erros comuns a evitar", bullets=[
    "Comparar 3+ grupos par a par com testes t (use ANOVA).",
    "Mudar para teste unilateral depois de olhar os dados.",
    "Reportar apenas o valor-p, sem tamanho de efeito ou intervalo de confiança.",
    "Tratar dados pareados como independentes (ou vice-versa).",
    "Confundir significância com relevância prática.",
    "Iniciar um estudo sem cálculo amostral.",
])

refs_slide()

prs.save(str(OUT))
print(f"Slides salvos: {OUT.name}  ({len(prs.slides._sldIdLst)} slides)")
